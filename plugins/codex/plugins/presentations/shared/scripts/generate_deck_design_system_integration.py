#!/usr/bin/env python3
"""Generate 023 deck-design-system integration mini-deck fixtures."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import generate_reference_calibrated_candidates as candidate_gen


SHARED = Path(__file__).resolve().parents[1]
REPO_ROOT = Path(__file__).resolve().parents[6]
COMPOSITION_INDEX = SHARED / "references" / "research_slide_composition_index.json"

TASK_KEY = "023_research_presentation_deck_design_system_integration"
PROFILE_ID = "research-deck-design-profile-v1"
SLIDE_W = 13.333
SLIDE_H = 7.5

INK = "171C26"
MUTED = "586274"
BG = "FCFCFA"
LINE = "BCC4D0"
ACCENT = "006970"
SECONDARY = "3F4578"
WARNING = "AE5E19"

FORBIDDEN_AUDIENCE_TERMS = [
    "RRL-",
    "SRC-",
    "candidate",
    "Reference retrieval",
    "EVIDENCE_MANIFEST",
    "Diagram contract",
    "QA",
    "repo path",
    "run ID",
    "implementation commit",
    "review target",
    "reference_faithful",
    "alternative_composition",
    "controlled_wildcard",
]


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def stable_json_sha(payload: Any) -> str:
    return hashlib.sha256(json.dumps(payload, sort_keys=True, separators=(",", ":")).encode("utf-8")).hexdigest()


def rel(path: Path) -> str:
    return str(path.resolve().relative_to(REPO_ROOT))


def records_by_id() -> dict[str, dict[str, Any]]:
    return {record["reference_id"]: record for record in json.loads(COMPOSITION_INDEX.read_text(encoding="utf-8"))["records"]}


def primary_region(record: dict[str, Any]) -> dict[str, Any]:
    return next(region for region in record["regions"] if region["region_id"] == record["primary_scientific_object_region_id"])


def bbox_to_inches(bbox: dict[str, float]) -> tuple[float, float, float, float]:
    return (
        float(bbox["x"]) * SLIDE_W,
        float(bbox["y"]) * SLIDE_H,
        float(bbox["w"]) * SLIDE_W,
        float(bbox["h"]) * SLIDE_H,
    )


def color(value: str) -> RGBColor:
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_text(slide: Any, text: str, bbox: dict[str, float], *, size: int, bold: bool = False, fill: str = INK, align: int | None = None) -> Any:
    x, y, w, h = bbox_to_inches(bbox)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)
    return shape


def add_shape(slide: Any, kind: Any, bbox: dict[str, float], *, line: str = LINE, fill: str | None = None, width: float = 1.0) -> Any:
    x, y, w, h = bbox_to_inches(bbox)
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(width)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(fill)
    else:
        shape.fill.background()
    return shape


def add_picture_contain(slide: Any, image_path: Path, bbox: dict[str, float]) -> Any:
    with Image.open(image_path) as img:
        iw, ih = img.size
    x, y, w, h = bbox_to_inches(bbox)
    scale = min(w / iw, h / ih)
    draw_w = iw * scale
    draw_h = ih * scale
    px = x + (w - draw_w) / 2
    py = y + (h - draw_h) / 2
    return slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(draw_w), height=Inches(draw_h))


def title_bbox() -> dict[str, float]:
    return {"x": 0.065, "y": 0.055, "w": 0.84, "h": 0.105}


def caption_bbox() -> dict[str, float]:
    return {"x": 0.075, "y": 0.865, "w": 0.76, "h": 0.055}


def source_primary_bbox(reference_id: str) -> dict[str, float]:
    return primary_region(records_by_id()[reference_id])["bbox"]


def deck_design_profile() -> dict[str, Any]:
    locked = {
        "fonts": {
            "primary": "Aptos",
            "secondary": "Aptos",
            "equation_asset_policy": "vector-or-high-resolution-raster-asset-inside-native-PPTX-frame",
        },
        "type_scale": {
            "title_pt": 27,
            "subtitle_pt": 15,
            "body_pt": 17,
            "caption_pt": 10,
            "annotation_pt": 13,
            "equation_min_height": 0.16,
        },
        "color_roles": {
            "background": f"#{BG}",
            "ink": f"#{INK}",
            "muted": f"#{MUTED}",
            "line": f"#{LINE}",
            "accent": f"#{ACCENT}",
            "secondary": f"#{SECONDARY}",
            "warning": f"#{WARNING}",
            "uncertainty": "#7A8798",
        },
        "spacing": {
            "outer_margin": 0.065,
            "object_gap": 0.026,
            "annotation_gap": 0.030,
            "panel_label_gap": 0.012,
        },
        "annotation": {
            "leader_color_role": "accent",
            "leader_width_pt": 2.0,
            "callout_rule": "left-rule plus direct leader to scientific object",
        },
        "chart": {
            "axis_color_role": "muted",
            "uncertainty_style": "thin muted interval or nominal reference line",
            "legend_position": "adjacent to plotted evidence",
        },
        "image_panel": {
            "label_position": "adjacent-above-panel",
            "legend_binding": "shared legend aligned with image evidence",
            "container_role": "semantic-panel-border-only",
        },
        "equation": {
            "background": "slide-canvas",
            "highlight_role": "warning",
            "leader_role": "accent",
            "decorative_card_used": False,
        },
        "caption": {
            "position": "near supporting evidence",
            "style": "short scientific qualifier",
            "internal_provenance_allowed": False,
        },
    }
    return {
        "schema": "RESEARCH_DECK_DESIGN_PROFILE_V1",
        "profile_id": PROFILE_ID,
        "provenance": {
            "task_key": TASK_KEY,
            "source_layers": ["019 composition records", "020 candidate search", "022 visual-finish primitives"],
            "audience": "research group meeting",
            "medium": "editable PPTX rendered through real presentation engine",
            "content_density": "compact specialist scientific update",
        },
        "locked_properties": locked,
        "locked_properties_sha256": stable_json_sha(locked),
        "page_local_properties": [
            "scientific_object_bbox",
            "layout_family",
            "panel_count",
            "primary_object_role",
            "annotation_target",
            "page_specific_accent_placement",
            "relative_positions_of_equation_image_plot_diagram",
        ],
    }


def base_request(request_name: str) -> dict[str, Any]:
    path = REPO_ROOT / "docs/audits/research_presentation_candidate_search/requests" / request_name
    return json.loads(path.read_text(encoding="utf-8"))


def slot_map(request: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {slot["slot_id"]: slot for slot in request["content_slots"]}


def add_slide_bg(slide: Any) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color(BG)


def add_footer(slide: Any, text: str) -> None:
    add_text(slide, text, caption_bbox(), size=10, fill=MUTED)


def add_equation_slide(prs: Presentation, profile: dict[str, Any], slide_spec: dict[str, Any]) -> dict[str, Any]:
    request = base_request("statistical_estimator_request.json")
    source = records_by_id()["RRL-028"]
    regions, transfers, family, reading_flow = candidate_gen.candidate_regions(request, "reference_faithful", source)
    render_regions = []
    for item in regions:
        copy = dict(item)
        if copy["region_id"] == "annotation":
            copy["bbox"] = {"x": 0.62, "y": 0.57, "w": 0.26, "h": 0.10}
        render_regions.append(copy)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    slots = slot_map(request)
    for region in render_regions:
        bbox = region["bbox"]
        slot = slots[region["content_slot_id"]]
        if region["role"] == "title":
            add_text(slide, slide_spec["title"], title_bbox(), size=27, bold=True)
        elif region["content_mode"] == "equation":
            add_picture_contain(slide, REPO_ROOT / slot["asset_path"], bbox)
            x, y, w, h = bbox_to_inches(bbox)
            line = slide.shapes.add_shape(MSO_SHAPE.LEFT_BRACE, Inches(x + w * 0.42), Inches(y + h * 0.72), Inches(w * 0.35), Inches(h * 0.19))
            line.line.color.rgb = color(WARNING)
            line.line.width = Pt(2.5)
            line.fill.background()
        elif region["role"] == "annotation":
            add_text(slide, slots["annotation"]["text"], bbox, size=13, fill=INK)
            x, y, w, h = bbox_to_inches(bbox)
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y - 0.08), Inches(6.9), Inches(3.95))
            connector.line.color.rgb = color(ACCENT)
            connector.line.width = Pt(2)
        elif region["role"] == "caption":
            add_footer(slide, "Synthetic clustered-data fixture; design profile is deck-wide, geometry remains page-local.")
    return slide_manifest(slide_spec, source, family, reading_flow, render_regions, transfers, profile)


def add_plot_slide(prs: Presentation, profile: dict[str, Any], slide_spec: dict[str, Any], *, reference_id: str, image_path: Path, annotation: str) -> dict[str, Any]:
    source = records_by_id()[reference_id]
    primary = source_primary_bbox(reference_id)
    title = title_bbox()
    plot_bbox = {
        "x": max(0.09, float(primary["x"]) - 0.02),
        "y": max(0.20, float(primary["y"])),
        "w": min(0.80, max(0.62, float(primary["w"]) + 0.08)),
        "h": min(0.58, max(0.42, float(primary["h"]) + 0.04)),
    }
    ann_bbox = {"x": plot_bbox["x"] + 0.03, "y": min(0.78, plot_bbox["y"] + plot_bbox["h"] + 0.025), "w": min(0.70, plot_bbox["w"] - 0.06), "h": 0.07}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_text(slide, slide_spec["title"], title, size=27, bold=True)
    add_picture_contain(slide, image_path, plot_bbox)
    add_text(slide, annotation, ann_bbox, size=13, fill=INK)
    add_footer(slide, slide_spec["caption"])
    regions = [
        region("title", "title", title, "title", "text"),
        region("plot", "primary_scientific_object", plot_bbox, "plot", "figure"),
        region("annotation", "annotation", ann_bbox, "annotation", "text"),
        region("caption", "caption", caption_bbox(), "caption", "caption"),
    ]
    transfers = [transfer(source, primary_region(source), "plot", plot_bbox, "scale")]
    return slide_manifest(slide_spec, source, source["layout_family"], source["reading_flow"], regions, transfers, profile)


def add_flow_slide(prs: Presentation, profile: dict[str, Any], slide_spec: dict[str, Any], *, reference_id: str, nodes: list[str], caption: str) -> dict[str, Any]:
    source = records_by_id()[reference_id]
    primary = source_primary_bbox(reference_id)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_text(slide, slide_spec["title"], title_bbox(), size=27, bold=True)
    flow_bbox = {
        "x": max(0.08, float(primary["x"])),
        "y": max(0.26, float(primary["y"])),
        "w": min(0.84, max(0.70, float(primary["w"]))),
        "h": min(0.30, max(0.20, float(primary["h"]))),
    }
    gap = 0.04
    node_w = (flow_bbox["w"] - gap * (len(nodes) - 1)) / len(nodes)
    regions = [region("title", "title", title_bbox(), "title", "text")]
    y_mid = flow_bbox["y"] + flow_bbox["h"] * 0.50
    start_x = flow_bbox["x"]
    end_x = flow_bbox["x"] + flow_bbox["w"]
    spine = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start_x * SLIDE_W), Inches(y_mid * SLIDE_H), Inches(end_x * SLIDE_W), Inches(y_mid * SLIDE_H))
    spine.line.color.rgb = color(ACCENT)
    spine.line.width = Pt(2.5)
    previous_right: tuple[float, float] | None = None
    for index, label in enumerate(nodes):
        bbox = {
            "x": flow_bbox["x"] + index * (node_w + gap),
            "y": flow_bbox["y"],
            "w": node_w,
            "h": flow_bbox["h"],
        }
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches((bbox["x"] + bbox["w"] / 2) * SLIDE_W - 0.08), Inches(y_mid * SLIDE_H - 0.08), Inches(0.16), Inches(0.16))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color(BG)
        marker.line.color.rgb = color(ACCENT)
        marker.line.width = Pt(2)
        add_text(slide, label, {"x": bbox["x"], "y": bbox["y"] + 0.02, "w": bbox["w"], "h": 0.08}, size=14, bold=True, fill=INK, align=PP_ALIGN.CENTER)
        add_text(slide, stage_microcopy(label), {"x": bbox["x"] + 0.008, "y": y_mid + 0.055, "w": bbox["w"] - 0.016, "h": 0.10}, size=10, fill=MUTED, align=PP_ALIGN.CENTER)
        regions.append(region(f"node_{index + 1}", "primary_scientific_object" if index == 1 else "secondary_scientific_object", bbox, f"node_{index + 1}", "diagram"))
        if previous_right:
            pass
        previous_right = ((bbox["x"] + bbox["w"]) * SLIDE_W, (bbox["y"] + bbox["h"] / 2) * SLIDE_H)
    add_footer(slide, caption)
    regions.append(region("caption", "caption", caption_bbox(), "caption", "caption"))
    transfers = [transfer(source, primary_region(source), "node_2", flow_bbox, "split")]
    return slide_manifest(slide_spec, source, source["layout_family"], source["reading_flow"], regions, transfers, profile)


def stage_microcopy(label: str) -> str:
    return {
        "DGP knobs": "ICC, centers, imbalance",
        "Clustered samples": "same estimand",
        "Interval procedures": "naive vs CR",
        "Coverage / width": "endpoint evidence",
        "Input image": "same case",
        "GT / prediction": "paired overlay",
        "Endpoint summary": "quantified error",
        "Failure review": "case-level readout",
    }.get(label, "")


def add_medical_slide(prs: Presentation, profile: dict[str, Any], slide_spec: dict[str, Any], *, strategy: str) -> dict[str, Any]:
    request = base_request("medical_image_comparison_request.json")
    source_id = "RRL-013" if strategy == "alternative_composition" else "RRL-022"
    source = records_by_id()[source_id]
    regions, transfers, family, reading_flow = candidate_gen.candidate_regions(request, strategy, source)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    slots = slot_map(request)
    for item in regions:
        slot = slots[item["content_slot_id"]]
        bbox = item["bbox"]
        if item["role"] == "title":
            add_text(slide, slide_spec["title"], title_bbox(), size=27, bold=True)
        elif item["content_mode"] == "medical_image":
            label_box = {"x": bbox["x"], "y": max(0.17, bbox["y"] - 0.045), "w": bbox["w"], "h": 0.035}
            add_text(slide, panel_label(slot["slot_id"]), label_box, size=10, bold=True, fill=MUTED)
            add_shape(slide, MSO_SHAPE.RECTANGLE, bbox, line=LINE, width=0.8)
            add_picture_contain(slide, REPO_ROOT / slot["asset_path"], bbox)
        elif item["role"] == "legend":
            add_legend(slide, bbox)
        elif item["role"] == "annotation":
            add_text(slide, slot["text"], bbox, size=13, fill=INK)
    add_footer(slide, "Synthetic public-safe image fixture; image treatment is integrated, clinical realism remains out of scope.")
    return slide_manifest(slide_spec, source, family, reading_flow, regions, transfers, profile)


def add_legend(slide: Any, bbox: dict[str, float]) -> None:
    labels = [("2D8C53", "overlap"), ("CA4A4A", "false positive"), ("4874C4", "false negative")]
    x, y, w, h = bbox_to_inches(bbox)
    cursor = x
    for hex_color, label in labels:
        swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cursor), Inches(y + 0.05), Inches(0.14), Inches(0.14))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color(hex_color)
        swatch.line.color.rgb = color(hex_color)
        add_text(slide, label, {"x": cursor / SLIDE_W + 0.015, "y": (y + 0.02) / SLIDE_H, "w": 0.13, "h": h / SLIDE_H}, size=9, fill=MUTED)
        cursor += 1.55


def panel_label(slot_id: str) -> str:
    return {
        "input_image": "Input",
        "overlay_image": "GT / prediction",
        "prediction_image": "Prediction",
        "error_image": "Error map",
    }.get(slot_id, slot_id.replace("_", " ").title())


def region(region_id: str, role: str, bbox: dict[str, float], slot_id: str, mode: str) -> dict[str, Any]:
    return {
        "region_id": region_id,
        "role": role,
        "bbox": {key: round(float(bbox[key]), 4) for key in ["x", "y", "w", "h"]},
        "content_slot_id": slot_id,
        "content_mode": mode,
    }


def transfer(source: dict[str, Any], source_item: dict[str, Any], candidate_region_id: str, candidate_bbox: dict[str, float], adaptation_type: str) -> dict[str, Any]:
    return {
        "source_reference_id": source["reference_id"],
        "source_region_id": source_item["region_id"],
        "source_role": source_item["role"],
        "source_bbox": source_item["bbox"],
        "candidate_region_id": candidate_region_id,
        "candidate_bbox": {key: round(float(candidate_bbox[key]), 4) for key in ["x", "y", "w", "h"]},
        "adaptation_type": adaptation_type,
        "adaptation_reason": "used source-derived normalized geometry while applying locked deck design profile",
    }


def slide_manifest(slide_spec: dict[str, Any], source: dict[str, Any], family: str, reading_flow: str, regions: list[dict[str, Any]], transfers: list[dict[str, Any]], profile: dict[str, Any]) -> dict[str, Any]:
    primary = [item for item in regions if item["role"] == "primary_scientific_object"]
    return {
        "slide_id": slide_spec["slide_id"],
        "title": slide_spec["title"],
        "page_function": slide_spec["page_function"],
        "deck_role": slide_spec["deck_role"],
        "source_reference_ids": [source["reference_id"]],
        "source_page_function": source["page_function"],
        "source_layout_family": source["layout_family"],
        "layout_family": family,
        "reading_flow": reading_flow,
        "primary_object_roles": sorted({item["content_mode"] for item in primary}),
        "primary_bboxes": [item["bbox"] for item in primary],
        "regions": regions,
        "geometry_transfer": transfers,
        "locked_profile_id": profile["profile_id"],
        "locked_properties_sha256": profile["locked_properties_sha256"],
        "page_local_geometry_preserved": True,
        "audience_text": [slide_spec["title"], slide_spec.get("caption", "")],
    }


def statistical_specs() -> list[dict[str, Any]]:
    return [
        {"slide_id": "stat_s01", "deck_role": "estimator", "page_function": "ESTIMATOR", "title": "Cluster-robust variance centers uncertainty at the group level", "caption": "Synthetic clustered-data fixture."},
        {"slide_id": "stat_s02", "deck_role": "simulation_design", "page_function": "EXPERIMENT_DESIGN", "title": "Simulation varies clustering before interval evaluation", "caption": "DGP knobs feed generated samples and interval procedures."},
        {"slide_id": "stat_s03", "deck_role": "main_result", "page_function": "RESULT_FIGURE", "title": "Cluster-robust intervals recover coverage, but small-G remains stressed", "caption": "Deterministic simulation; nominal 0.95 reference line shown in the plotted evidence."},
        {"slide_id": "stat_s04", "deck_role": "negative_result", "page_function": "NEGATIVE_RESULT", "title": "Small center counts remain anti-conservative", "caption": "Negative result motivates CR2 or wild cluster bootstrap as the next discriminating experiment."},
    ]


def medical_specs() -> list[dict[str, Any]]:
    return [
        {"slide_id": "med_s01", "deck_role": "image_target", "page_function": "MEDICAL_IMAGE_COMPARISON", "title": "Overlay errors localize the small-lesion failure mode", "caption": "Synthetic public-safe image fixture."},
        {"slide_id": "med_s02", "deck_role": "method_path", "page_function": "EXPERIMENT_DESIGN", "title": "The evaluation path keeps image evidence and endpoints paired", "caption": "Input, mask, prediction, and endpoint summaries are generated from the same synthetic case."},
        {"slide_id": "med_s03", "deck_role": "quantitative_result", "page_function": "RESULT_FIGURE", "title": "Lesion-size recall exposes the endpoint-level weakness", "caption": "Deterministic fixture result; not clinical validation."},
        {"slide_id": "med_s04", "deck_role": "same_case_error", "page_function": "MEDICAL_IMAGE_COMPARISON", "title": "Same-case error maps show where the endpoint is lost", "caption": "Synthetic case-level comparison keeps GT, prediction, and error map tied together."},
    ]


def build_deck(deck_key: str, specs: list[dict[str, Any]], profile: dict[str, Any], out_dir: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    slides: list[dict[str, Any]] = []
    stat_assets = REPO_ROOT / "tests/fixtures/presentations/statistical_method_group_meeting/visual_review_packet_source/assets"
    med_assets = REPO_ROOT / "tests/fixtures/presentations/medical_imaging_group_meeting/visual_review_packet_source/assets"
    for spec in specs:
        if spec["slide_id"] == "stat_s01":
            slides.append(add_equation_slide(prs, profile, spec))
        elif spec["slide_id"] == "stat_s02":
            slides.append(add_flow_slide(prs, profile, spec, reference_id="RRL-019", nodes=["DGP knobs", "Clustered samples", "Interval procedures", "Coverage / width"], caption=spec["caption"]))
        elif spec["slide_id"] == "stat_s03":
            slides.append(add_plot_slide(prs, profile, spec, reference_id="RRL-030", image_path=stat_assets / "coverage_by_icc.png", annotation="Coverage improves after center-level aggregation, with visible stress at low center counts."))
        elif spec["slide_id"] == "stat_s04":
            slides.append(add_plot_slide(prs, profile, spec, reference_id="RRL-025", image_path=stat_assets / "small_g_negative_result.png", annotation="The failure region is the small-G corner; the next experiment should target finite-sample correction."))
        elif spec["slide_id"] == "med_s01":
            slides.append(add_medical_slide(prs, profile, spec, strategy="controlled_wildcard"))
        elif spec["slide_id"] == "med_s02":
            slides.append(add_flow_slide(prs, profile, spec, reference_id="RRL-019", nodes=["Input image", "GT / prediction", "Endpoint summary", "Failure review"], caption=spec["caption"]))
        elif spec["slide_id"] == "med_s03":
            slides.append(add_plot_slide(prs, profile, spec, reference_id="RRL-023", image_path=med_assets / "slide5_lesion_size_recall.png", annotation="Recall drops most sharply for the smallest lesions in this deterministic fixture."))
        elif spec["slide_id"] == "med_s04":
            slides.append(add_medical_slide(prs, profile, spec, strategy="alternative_composition"))
        else:
            raise RuntimeError(f"unknown slide spec {spec['slide_id']}")

    out_dir.mkdir(parents=True, exist_ok=True)
    deck_plan_path = out_dir / "DECK_PLAN.json"
    deck_plan_path.write_text(json.dumps({
        "schema": "RESEARCH_DECK_DESIGN_SYSTEM_FIXTURE_PLAN_V1",
        "task_key": TASK_KEY,
        "deck_key": deck_key,
        "fixture_boundary": "engineering integration fixture, not final design-quality holdout",
        "deck_design_profile": rel(out_dir.parent / "deck_design_profile.json"),
        "slides": specs,
    }, indent=2), encoding="utf-8")
    pptx_path = out_dir / f"{deck_key}.pptx"
    prs.save(pptx_path)
    render_status = render_pptx(pptx_path, out_dir)
    manifest = {
        "schema": "RESEARCH_DECK_DESIGN_SYSTEM_INTEGRATION_MANIFEST_V1",
        "task_key": TASK_KEY,
        "deck_key": deck_key,
        "fixture_boundary": "engineering integration fixture, not final design-quality holdout",
        "deck_design_profile": rel(out_dir.parent / "deck_design_profile.json"),
        "deck_plan": rel(deck_plan_path),
        "locked_properties_sha256": profile["locked_properties_sha256"],
        "pptx": rel(pptx_path),
        "pptx_sha256": sha256(pptx_path),
        "editable_slide_count": editable_slide_count(pptx_path),
        "slides": slides,
        "major_composition_families": sorted({slide["layout_family"].split("-reordered")[0].split("-focus")[0] for slide in slides}),
        "render_status": render_status,
        "mechanical_qa": {},
    }
    manifest["mechanical_qa"] = mechanical_qa(manifest, pptx_path)
    manifest_path = out_dir / "SOURCE_GENERATED_IDENTITY_MANIFEST.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    (out_dir / "MECHANICAL_VISUAL_REVIEW.json").write_text(json.dumps(manifest["mechanical_qa"], indent=2), encoding="utf-8")
    return {"manifest": rel(manifest_path), "pptx": rel(pptx_path), "render_status": rel(out_dir / "RENDER_STATUS.json")}


def build_review_pack(decks: list[dict[str, Any]], pack_path: Path) -> dict[str, Any]:
    pdfunite = shutil.which("pdfunite")
    pdfs = []
    for deck in decks:
        render_status = json.loads((REPO_ROOT / deck["render_status"]).read_text(encoding="utf-8"))
        pdf = REPO_ROOT / render_status.get("pptx_to_pdf", "")
        if not pdf.exists():
            raise RuntimeError(f"rendered deck PDF missing for review pack: {pdf}")
        pdfs.append(pdf)
    if not pdfunite:
        raise RuntimeError("pdfunite not found; cannot build combined review pack")
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    result = subprocess.run([pdfunite, *map(str, pdfs), str(pack_path)], check=False, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"pdfunite failed: {result.stderr[-2000:]}")
    return {
        "path": rel(pack_path),
        "source_pdfs": [rel(pdf) for pdf in pdfs],
        "sha256": sha256(pack_path),
    }


def editable_slide_count(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])


def find_renderer() -> str | None:
    if os.environ.get("RESEARCH_PPTX_RENDERER"):
        return os.environ["RESEARCH_PPTX_RENDERER"]
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    bundled = REPO_ROOT / ".cache/tools/squashfs-root/opt/libreoffice25.8/program/soffice"
    if bundled.exists():
        return str(bundled)
    return None


def render_pptx(pptx_path: Path, out_dir: Path) -> dict[str, Any]:
    renderer = find_renderer()
    pdf_dir = out_dir / "pdf"
    png_dir = out_dir / "rendered"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)
    if not renderer:
        status = {"status": "BLOCKED_REAL_PPTX_RENDER", "reason": "no soffice/libreoffice/RESEARCH_PPTX_RENDERER found", "png_count": 0}
        (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    with tempfile.TemporaryDirectory() as tmp_profile:
        cmd = [
            renderer,
            "--headless",
            f"-env:UserInstallation=file://{tmp_profile}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(pptx_path.resolve()),
        ]
        convert = subprocess.run(cmd, check=False, capture_output=True, text=True, timeout=120)
    if convert.returncode != 0:
        status = {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pptx to pdf conversion failed", "stderr": convert.stderr[-2000:], "png_count": 0}
        (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    pdf_candidates = sorted(pdf_dir.glob("*.pdf"))
    if not pdf_candidates:
        status = {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pptx to pdf conversion produced no PDF", "png_count": 0}
        (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        status = {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pdftoppm not found", "png_count": 0}
        (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    prefix = png_dir / "page"
    ppm = subprocess.run([pdftoppm, "-png", "-r", "160", str(pdf_candidates[0]), str(prefix)], check=False, capture_output=True, text=True, timeout=120)
    if ppm.returncode != 0:
        status = {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pdftoppm failed", "stderr": ppm.stderr[-2000:], "png_count": 0}
        (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    pages = sorted(png_dir.glob("page-*.png"))
    renamed = []
    for index, page in enumerate(pages, start=1):
        target = png_dir / f"slide-{index}.png"
        if target.exists():
            target.unlink()
        page.rename(target)
        renamed.append(target)
    status = {
        "status": "ok",
        "renderer": renderer,
        "pptx_to_pdf": rel(pdf_candidates[0]),
        "png_count": len(renamed),
        "rendered_png": [{"path": rel(path), "sha256": sha256(path)} for path in renamed],
    }
    (out_dir / "RENDER_STATUS.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
    return status


def pptx_text(pptx_path: Path) -> str:
    with ZipFile(pptx_path) as zf:
        chunks = []
        for name in sorted(zf.namelist()):
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                chunks.append(zf.read(name).decode("utf-8", errors="ignore"))
        return "\n".join(chunks)


def mechanical_qa(manifest: dict[str, Any], pptx_path: Path) -> dict[str, Any]:
    text = pptx_text(pptx_path)
    forbidden_hits = [term for term in FORBIDDEN_AUDIENCE_TERMS if term in text]
    families = manifest["major_composition_families"]
    profile_shas = {slide["locked_properties_sha256"] for slide in manifest["slides"]}
    slide_count = editable_slide_count(pptx_path)
    render = manifest["render_status"]
    checks = {
        "editable_pptx": slide_count == len(manifest["slides"]),
        "real_render": render.get("status") == "ok" and render.get("png_count") == len(manifest["slides"]),
        "deck_wide_profile_locked": len(profile_shas) == 1,
        "composition_diversity": len(families) >= 3,
        "audience_meta_leakage": not forbidden_hits,
        "fixture_boundary_recorded": manifest["fixture_boundary"].startswith("engineering integration fixture"),
    }
    return {
        "schema": "RESEARCH_DECK_DESIGN_SYSTEM_MECHANICAL_QA_V1",
        "status": "MECHANICAL_PASS" if all(checks.values()) else "MECHANICAL_REVISE",
        "checks": checks,
        "forbidden_audience_hits": forbidden_hits,
        "rendered_png_count": render.get("png_count", 0),
        "pptx_object_summary": {
            "editable_slide_count": slide_count,
            "major_composition_families": families,
            "locked_properties_sha256": sorted(profile_shas),
        },
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out-dir", default=f"results/{TASK_KEY}/generated")
    args = parser.parse_args()
    out_root = REPO_ROOT / args.out_dir
    if out_root.exists():
        shutil.rmtree(out_root)
    out_root.mkdir(parents=True)
    profile = deck_design_profile()
    profile_path = out_root / "deck_design_profile.json"
    profile_path.write_text(json.dumps(profile, indent=2), encoding="utf-8")
    decks = [
        build_deck("statistical_design_system_fixture", statistical_specs(), profile, out_root / "statistical_design_system_fixture"),
        build_deck("medical_design_system_fixture", medical_specs(), profile, out_root / "medical_design_system_fixture"),
    ]
    review_pack = build_review_pack(decks, out_root.parent / "REVIEW_PACK.pdf")
    outputs = {
        "schema": "RESEARCH_DECK_DESIGN_SYSTEM_INTEGRATION_OUTPUTS_V1",
        "task_key": TASK_KEY,
        "deck_design_profile": rel(profile_path),
        "decks": decks,
        "review_pack_pdf": review_pack,
    }
    outputs_path = out_root / "OUTPUTS.json"
    outputs_path.write_text(json.dumps(outputs, indent=2), encoding="utf-8")
    print(json.dumps({"outputs": rel(outputs_path), "decks": outputs["decks"]}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
