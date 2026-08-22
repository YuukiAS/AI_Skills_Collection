#!/usr/bin/env python3
"""Generate a five-slide statistical method group-meeting benchmark deck."""

from __future__ import annotations

import argparse
import csv
import json
import math
import os
import random
import shutil
import subprocess
import sys
from pathlib import Path
from zipfile import ZipFile

os.environ.setdefault("MPLCONFIGDIR", "/tmp/ai-skills-matplotlib")
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


W, H = 13.333, 7.5
REPO_ROOT = Path(__file__).resolve().parents[4]
REFERENCE_ROOT = REPO_ROOT / "skills" / "tools" / "documents-media" / "presentations" / "shared" / "references"
REFERENCE_INDEX = REFERENCE_ROOT / "research_slide_reference_index.csv"
REFERENCE_MANIFEST = REFERENCE_ROOT / "reference_sources_manifest.json"
SEED = 20260822
REPLICATES = 400
CENTER_GRID = [8, 20, 50]
ICC_GRID = [0.0, 0.1, 0.3, 0.5]
IMBALANCE_GRID = ["balanced", "imbalanced"]
BETA1 = 0.25
P = {
    "bg": "F1F2F4",
    "ink": "17202A",
    "muted": "606977",
    "line": "C6CCD6",
    "purple": "4F1F68",
    "teal": "0F766E",
    "blue": "1F4E79",
    "gold": "9A6A16",
    "red": "A33A34",
    "green": "2A7C54",
    "white": "FFFFFF",
    "soft_teal": "E1F4F1",
    "soft_blue": "E5EEF8",
    "soft_gold": "FBF1D6",
    "soft_red": "F8E3E1",
}

REFERENCE_QUERIES = {
    "STATISTICAL_MODEL": {
        "intent": "Explain estimand, DGP, cluster correlation, and why unbiased point estimates can still have wrong interval coverage.",
        "page_functions": ["STATISTICAL_MODEL", "BAYESIAN_MODEL", "ASSUMPTION", "MODEL_CHECK"],
        "scientific_domain": ["statistics", "biostatistics"],
        "statistical_subdomain": ["Bayesian workflow", "Bayesian data analysis", "survey nonresponse / MRP", "Bayesian clinical research"],
        "evidence_types": ["model objective", "Bayesian model", "assumption / challenge", "modeling workflow", "confidence interval principle"],
        "organization_lesson": "Make the estimand and data-generating mechanism the main object before discussing results.",
    },
    "ESTIMATOR": {
        "intent": "Show how a variance estimator changes when score contributions are aggregated by center.",
        "page_functions": ["ESTIMATOR", "CONFIDENCE_INTERVAL", "STATISTICAL_MODEL", "ASSUMPTION"],
        "scientific_domain": ["statistics", "biostatistics"],
        "statistical_subdomain": ["hybrid resource-bound analysis", "survey nonresponse / MRP", "Bayesian clinical research", "Bayesian workflow"],
        "evidence_types": ["estimator formula", "estimator pipeline", "uncertainty interval", "confidence interval principle", "metric definition"],
        "organization_lesson": "Put the formula next to definitions of the grouped objects it operates on.",
    },
    "SIMULATION_DESIGN": {
        "intent": "Connect DGP knobs, generated centers, interval methods, replicates, and coverage endpoints in one experiment design.",
        "page_functions": ["EXPERIMENT_DESIGN", "METHOD_DIAGRAM", "SIMULATION", "NEXT_EXPERIMENT", "ESTIMATOR"],
        "scientific_domain": ["statistics", "biostatistics", "medical imaging"],
        "statistical_subdomain": ["Bayesian workflow", "hybrid resource-bound analysis", "survey nonresponse / MRP", "lesion segmentation"],
        "evidence_types": ["task overview", "method mechanism", "simulation comparison", "planned evidence", "estimator pipeline"],
        "organization_lesson": "Use one left-to-right experimental path with explicit endpoints and visible comparator branches.",
    },
    "RESULT_FIGURE": {
        "intent": "Show coverage against ICC with nominal target line, method comparison, and Monte Carlo uncertainty.",
        "page_functions": ["RESULT_FIGURE", "CONFIDENCE_INTERVAL", "SIMULATION", "REAL_DATA_APPLICATION", "SENSITIVITY_ANALYSIS"],
        "scientific_domain": ["statistics", "biostatistics"],
        "statistical_subdomain": ["Bayesian workflow", "survey nonresponse / MRP", "Bayesian inference design", "Bayesian clinical research"],
        "evidence_types": ["quantitative plot", "uncertainty interval", "simulation comparison", "time-series interval", "forest plot"],
        "organization_lesson": "Let the plot dominate and keep uncertainty and interpretation adjacent to the graph.",
    },
    "NEGATIVE_RESULT": {
        "intent": "Show the stress regime where the current interval still fails and define the next discriminating experiment.",
        "page_functions": ["NEGATIVE_RESULT", "FINITE_SAMPLE", "MODEL_CHECK", "NEXT_EXPERIMENT", "SENSITIVITY_ANALYSIS"],
        "scientific_domain": ["statistics", "biostatistics"],
        "statistical_subdomain": ["Bayesian workflow", "Bayesian clinical research", "hybrid resource-bound analysis", "Bayesian inference design"],
        "evidence_types": ["negative/fix title", "limitation analysis", "posterior predictive check", "planned evidence", "sensitivity analysis"],
        "organization_lesson": "Keep the failure visible and separate completed evidence from the next planned method.",
    },
}

TIER_PRIORITY = {
    "PRIMARY_RESEARCH_PRESENTATION": 0,
    "SECONDARY_TEACHING_REFERENCE": 1,
    "PRESENTATION_GUIDANCE": 2,
    "CANDIDATE_BACKLOG": 3,
}


def rgb(name: str) -> RGBColor:
    value = P[name]
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def hex_color(name: str) -> str:
    return "#" + P[name]


def inch(value: float):
    return Inches(value)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float = 13, color: str = "ink", bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def rect(slide, text: str, x: float, y: float, w: float, h: float, fill: str = "white", color: str = "ink", size: float = 12, bold: bool = False, line: str = "line"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    add_text(slide, text, x + 0.06, y + 0.05, w - 0.12, h - 0.1, size, color, bold, PP_ALIGN.CENTER)
    return shape


def add_rule(slide, x1: float, y1: float, x2: float, y2: float, color: str = "line", width: float = 1.1, dash: bool = False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dash:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return connector


def arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = "purple", width: float = 1.7):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    ln = connector.line._get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    ln.append(tail)
    return connector


def header(slide, number: int, title: str, message: str):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0), inch(0), inch(W), inch(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb("purple")
    band.line.fill.background()
    add_text(slide, title, 0.55, 0.36, 9.9, 0.4, 20.5, "ink", True)
    add_text(slide, f"{number:02d}/05", 11.6, 0.42, 0.95, 0.22, 9, "muted", False, PP_ALIGN.RIGHT)
    add_text(slide, message, 0.65, 0.88, 11.7, 0.34, 12.2, "purple", True)


def load_font(size: int):
    for path in ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/dejavu/DejaVuSans.ttf"]:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def render_math_asset(path: Path, expression: str, fontsize: int = 36, color: str = "ink") -> dict[str, object]:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(1, 1), dpi=360)
    fig.patch.set_alpha(0)
    text = fig.text(0.0, 0.5, expression, fontsize=fontsize, color=hex_color(color), va="center", ha="left")
    fig.canvas.draw()
    bbox = text.get_window_extent(renderer=fig.canvas.get_renderer()).expanded(1.08, 1.35)
    fig.savefig(path, transparent=True, bbox_inches=bbox.transformed(fig.dpi_scale_trans.inverted()), pad_inches=0.04)
    plt.close(fig)
    with Image.open(path) as image:
        width, height = image.size
    return {
        "path": str(path),
        "source": expression,
        "format": "matplotlib_mathtext_png_transparent",
        "dpi": 360,
        "pixel_width": width,
        "pixel_height": height,
    }


def add_math(slide, assets: Path, manifest: dict, key: str, expression: str, x: float, y: float, w: float, fontsize: int = 36, color: str = "ink"):
    asset_path = assets / f"{key}.png"
    manifest.setdefault("math_assets", {})[key] = render_math_asset(asset_path, expression, fontsize=fontsize, color=color)
    return slide.shapes.add_picture(str(asset_path), inch(x), inch(y), width=inch(w))


def tokens(value: str) -> set[str]:
    return {part.lower() for part in value.replace("/", " ").replace("-", " ").replace("_", " ").split() if len(part) > 2}


def load_reference_rows() -> list[dict[str, str]]:
    manifest = json.loads(REFERENCE_MANIFEST.read_text(encoding="utf-8"))
    source_tiers = {item["source_id"]: item["source_tier"] for item in manifest["candidate_sources"]}
    source_domains = {item["source_id"]: item["domain_family"] for item in manifest["candidate_sources"]}
    source_subdomains = {item["source_id"]: item["statistical_subdomain"] for item in manifest["candidate_sources"]}
    rows = []
    with REFERENCE_INDEX.open(encoding="utf-8", newline="") as fh:
        for row in csv.DictReader(fh):
            if row.get("verification_status") != "inspected":
                continue
            if not (row.get("source_file_sha256") and row.get("rendered_page_sha256")):
                continue
            row["source_tier"] = source_tiers.get(row["source_id"], "")
            row["domain_family"] = source_domains.get(row["source_id"], "")
            row["statistical_subdomain"] = source_subdomains.get(row["source_id"], "")
            rows.append(row)
    return rows


def score_reference(row: dict[str, str], query: dict[str, object]) -> tuple[int, list[str]]:
    score = 0
    reasons: list[str] = []
    if row["page_function"] in query["page_functions"]:
        score += 10
        reasons.append(f"page_function={row['page_function']} matches query")
    if row["evidence_type"] in query["evidence_types"]:
        score += 6
        reasons.append(f"evidence_type={row['evidence_type']} matches query")
    if row["domain_family"] in query["scientific_domain"]:
        score += 4
        reasons.append(f"domain={row['domain_family']} matches query")
    if row["statistical_subdomain"] in query["statistical_subdomain"]:
        score += 4
        reasons.append(f"subdomain={row['statistical_subdomain']} matches query")
    query_terms = tokens(" ".join([str(query["intent"]), " ".join(query["page_functions"]), " ".join(query["evidence_types"])]))
    row_terms = tokens(" ".join([row["scientific_object"], row["why_this_specific_page_works"], row["what_to_learn"], row["short_page_specific_observation"]]))
    overlap = sorted(query_terms & row_terms)
    if overlap:
        score += min(5, len(overlap))
        reasons.append("semantic tokens overlap: " + ", ".join(overlap[:6]))
    tier_bonus = max(0, 3 - TIER_PRIORITY.get(row["source_tier"], 3))
    score += tier_bonus
    reasons.append(f"source_tier={row['source_tier']} priority_bonus={tier_bonus}")
    return score, reasons


def retrieve_references(archetype: str) -> dict[str, object]:
    query = REFERENCE_QUERIES[archetype]
    scored = []
    for row in load_reference_rows():
        score, reasons = score_reference(row, query)
        if score > 0:
            scored.append((score, TIER_PRIORITY.get(row["source_tier"], 3), row["reference_id"], row, reasons))
    scored.sort(key=lambda item: (-item[0], item[1], item[2]))
    candidate_rows = scored[:8]
    selected_rows = candidate_rows[:5]
    if len(selected_rows) < 2:
        raise RuntimeError(f"reference retrieval for {archetype} found fewer than two inspected references")
    return {
        "query": query,
        "candidate_ids": [row["reference_id"] for _, _, _, row, _ in candidate_rows],
        "selected_ids": [row["reference_id"] for _, _, _, row, _ in selected_rows],
        "source_tiers": {row["reference_id"]: row["source_tier"] for _, _, _, row, _ in selected_rows},
        "ranking_relevance_reason": {row["reference_id"]: "; ".join(reasons) for _, _, _, row, reasons in selected_rows},
        "selected_page_lessons": {
            row["reference_id"]: {
                "primary_object": row["scientific_object"],
                "page_function": row["page_function"],
                "evidence_type": row["evidence_type"],
                "information_density": row["approximate_figure_text_ratio"],
                "annotation_style": row["short_page_specific_observation"],
                "useful_lesson": row["what_to_learn"],
                "not_copied": row["what_not_to_copy"],
            }
            for _, _, _, row, _ in selected_rows
        },
        "organization_lesson": query["organization_lesson"],
        "what_was_not_copied": "No full-slide screenshots, source images, institutional styling, public figures, private clinical data, or source-specific visual identity were copied.",
    }


def reference_footer(refs: list[str], purpose: str) -> str:
    return f"Reference retrieval: {purpose} selected inspected pages {', '.join(refs)}; trace in EVIDENCE_MANIFEST; style not copied."


REFERENCE_DESIGN_DECISIONS = {
    "STATISTICAL_MODEL": {
        "adopted_design_decisions": [
            "Use the DGP formula as the dominant object rather than prose cards.",
            "Ground notation adjacent to the displayed formula with short semantic labels.",
            "Use containment and spacing for additive components instead of causal arrows.",
        ],
        "deliberately_not_adopted": [
            "Do not copy any source deck styling or full-slide imagery.",
            "Do not use reference IDs or retrieval traces as slide footer content.",
        ],
    },
    "ESTIMATOR": {
        "adopted_design_decisions": [
            "Make the sandwich covariance formula the largest object on the page.",
            "Annotate the middle score-aggregation term directly beside the formula.",
            "Keep the naive comparator compact so it supports rather than competes with the main formula.",
        ],
        "deliberately_not_adopted": [
            "Do not turn the variance comparison into three pastel explanation cards.",
            "Do not show source-like ASCII formula strings to the audience.",
        ],
    },
    "SIMULATION_DESIGN": {
        "adopted_design_decisions": [
            "Keep one left-to-right experimental path from DGP knobs to endpoints.",
            "Use only scientific operations as nodes and keep branch labels short.",
            "Use a restrained neutral layout with one accent color for the cluster-robust branch.",
        ],
        "deliberately_not_adopted": [
            "Do not show diagram QA contracts or reading instructions on the slide.",
            "Do not use decorative workflow boxes when a node is not a real scientific operation.",
        ],
    },
    "RESULT_FIGURE": {
        "adopted_design_decisions": [
            "Let the coverage plot occupy most of the slide area.",
            "Place uncertainty and the nominal target directly in the figure grammar.",
            "Use direct annotation at the small-G stress point instead of a separate reading card.",
        ],
        "deliberately_not_adopted": [
            "Do not move interpretation into pastel side cards.",
            "Do not describe the result as a general theorem.",
        ],
    },
    "NEGATIVE_RESULT": {
        "adopted_design_decisions": [
            "Use the same quantitative evidence grammar as the main result page.",
            "Make the undercoverage point visible and annotated in the chart.",
            "Keep the planned CR2 / wild-bootstrap experiment as a small secondary statement.",
        ],
        "deliberately_not_adopted": [
            "Do not create separate negative-result, mechanism, and next-experiment cards.",
            "Do not make planned methods look like completed evidence.",
        ],
    },
}


def reference_design_audit(archetype: str, retrieval: dict[str, object]) -> dict[str, object]:
    decisions = REFERENCE_DESIGN_DECISIONS[archetype]
    return {
        "selected_reference_ids": retrieval["selected_ids"],
        "page_specific_lessons": retrieval["selected_page_lessons"],
        "adopted_design_decisions": decisions["adopted_design_decisions"],
        "deliberately_not_adopted": decisions["deliberately_not_adopted"],
        "current_slide_application": retrieval["organization_lesson"],
    }


def invert_2x2(a: float, b: float, c: float, d: float) -> tuple[tuple[float, float], tuple[float, float]]:
    det = a * d - b * c
    if abs(det) < 1e-12:
        raise RuntimeError("singular two-column OLS design")
    return ((d / det, -b / det), (-c / det, a / det))


def simulate_cell(g_count: int, rho: float, imbalance: str, replicates: int, seed: int) -> dict[str, object]:
    rng = random.Random(seed + g_count * 1000 + int(rho * 100) + (17 if imbalance == "imbalanced" else 0))
    tau2 = rho / (1 - rho) if rho > 0 else 0.0
    naive_hits = 0
    cluster_hits = 0
    bias_sum = 0.0
    naive_width_sum = 0.0
    cluster_width_sum = 0.0
    for _ in range(replicates):
        data: list[tuple[int, int, float]] = []
        clusters: list[list[tuple[int, float]]] = []
        for g in range(g_count):
            if imbalance == "imbalanced":
                n = [8, 10, 14, 22, 42, 55, 16, 30][g % 8]
                p_treat = [0.18, 0.30, 0.45, 0.65, 0.80][g % 5]
            else:
                n = 24
                p_treat = 0.5
            u_g = rng.gauss(0.0, math.sqrt(tau2)) if tau2 else 0.0
            cluster_rows: list[tuple[int, float]] = []
            for _i in range(n):
                treatment = 1 if rng.random() < p_treat else 0
                y = BETA1 * treatment + u_g + rng.gauss(0.0, 1.0)
                data.append((g, treatment, y))
                cluster_rows.append((treatment, y))
            clusters.append(cluster_rows)
        n_total = len(data)
        sx = sum(x for _, x, _ in data)
        sy = sum(y for _, _, y in data)
        sxy = sum(x * y for _, x, y in data)
        inv = invert_2x2(n_total, sx, sx, sx)
        b0 = inv[0][0] * sy + inv[0][1] * sxy
        b1 = inv[1][0] * sy + inv[1][1] * sxy
        residuals = [(g, x, y - b0 - b1 * x) for g, x, y in data]
        rss = sum(e * e for _, _, e in residuals)
        sigma2_hat = rss / (n_total - 2)
        naive_se = math.sqrt(max(0.0, sigma2_hat * inv[1][1]))
        middle00 = middle01 = middle11 = 0.0
        offset = 0
        for rows in clusters:
            score0 = score1 = 0.0
            for x, _y in rows:
                _g, _x, e = residuals[offset]
                offset += 1
                score0 += e
                score1 += x * e
            middle00 += score0 * score0
            middle01 += score0 * score1
            middle11 += score1 * score1
        v11 = inv[1][0] * (middle00 * inv[0][1] + middle01 * inv[1][1]) + inv[1][1] * (middle01 * inv[0][1] + middle11 * inv[1][1])
        correction = g_count / (g_count - 1) * (n_total - 1) / (n_total - 2)
        cluster_se = math.sqrt(max(0.0, v11 * correction))
        naive_lo, naive_hi = b1 - 1.96 * naive_se, b1 + 1.96 * naive_se
        cluster_lo, cluster_hi = b1 - 1.96 * cluster_se, b1 + 1.96 * cluster_se
        naive_hits += naive_lo <= BETA1 <= naive_hi
        cluster_hits += cluster_lo <= BETA1 <= cluster_hi
        bias_sum += b1 - BETA1
        naive_width_sum += naive_hi - naive_lo
        cluster_width_sum += cluster_hi - cluster_lo
    naive = naive_hits / replicates
    cluster = cluster_hits / replicates
    return {
        "G": g_count,
        "rho": rho,
        "imbalance": imbalance,
        "replicates": replicates,
        "methods": {
            "naive_iid_ols_z": {
                "coverage": round(naive, 4),
                "mc_se": round(math.sqrt(naive * (1 - naive) / replicates), 4),
                "mean_width": round(naive_width_sum / replicates, 4),
            },
            "cluster_robust_z": {
                "coverage": round(cluster, 4),
                "mc_se": round(math.sqrt(cluster * (1 - cluster) / replicates), 4),
                "mean_width": round(cluster_width_sum / replicates, 4),
            },
        },
        "bias": round(bias_sum / replicates, 4),
    }


def run_simulation(out_dir: Path) -> dict[str, object]:
    rows = []
    for imbalance in IMBALANCE_GRID:
        for g_count in CENTER_GRID:
            for rho in ICC_GRID:
                rows.append(simulate_cell(g_count, rho, imbalance, REPLICATES, SEED))
    stress = next(row for row in rows if row["G"] == 8 and row["rho"] == 0.5 and row["imbalance"] == "imbalanced")
    summary = {
        "seed": SEED,
        "replicates_per_cell": REPLICATES,
        "dgp": {
            "formula": "Y_ij = beta_0 + beta_1 T_ij + u_j + epsilon_ij",
            "beta_1": BETA1,
            "u_j": "N(0, tau^2)",
            "epsilon_ij": "N(0, sigma^2)",
            "icc": "rho = tau^2 / (tau^2 + sigma^2)",
        },
        "grid": {
            "center_count": CENTER_GRID,
            "icc": ICC_GRID,
            "imbalance": IMBALANCE_GRID,
        },
        "methods": ["naive_iid_ols_z", "cluster_robust_z"],
        "endpoints": ["95% interval coverage", "mean bias", "mean interval width"],
        "nominal_coverage": 0.95,
        "rows": rows,
        "negative_result": {
            "condition": "G=8, rho=0.5, imbalanced cluster sizes and treatment shares",
            "cluster_robust_coverage": stress["methods"]["cluster_robust_z"]["coverage"],
            "naive_iid_coverage": stress["methods"]["naive_iid_ols_z"]["coverage"],
            "claim": "small-G/high-ICC stress remains under nominal for cluster-robust z intervals",
            "next_experiment": "planned: compare CR2 small-sample correction and wild cluster bootstrap under the same DGP grid",
        },
    }
    sim_dir = out_dir / "simulation"
    sim_dir.mkdir(parents=True, exist_ok=True)
    (sim_dir / "simulation_summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    with (sim_dir / "simulation_summary.csv").open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=["G", "rho", "imbalance", "method", "coverage", "mc_se", "mean_width", "bias"])
        writer.writeheader()
        for row in rows:
            for method, values in row["methods"].items():
                writer.writerow({
                    "G": row["G"],
                    "rho": row["rho"],
                    "imbalance": row["imbalance"],
                    "method": method,
                    "coverage": values["coverage"],
                    "mc_se": values["mc_se"],
                    "mean_width": values["mean_width"],
                    "bias": row["bias"],
                })
    return summary


def draw_coverage_plot(path: Path, summary: dict[str, object]) -> None:
    img = Image.new("RGB", (1480, 700), "white")
    draw = ImageDraw.Draw(img)
    font = load_font(21)
    small = load_font(16)
    tiny = load_font(14)
    title = load_font(27)
    draw.text((54, 24), "Imbalanced clusters: 95% interval coverage by ICC", fill=hex_color("ink"), font=title)
    panels = [(8, 92), (20, 555), (50, 1018)]
    rows = summary["rows"]
    method_colors = {"naive_iid_ols_z": hex_color("red"), "cluster_robust_z": hex_color("teal")}
    for g_count, x0 in panels:
        y0 = 575
        plot_w = 360
        plot_h = 410
        draw.rectangle((x0, 122, x0 + plot_w, y0), outline=hex_color("line"), width=2)
        draw.text((x0 + 135, 88), f"G={g_count}", fill=hex_color("ink"), font=font)
        nominal_y = y0 - int((0.95 - 0.45) / 0.55 * plot_h)
        draw.line((x0, nominal_y, x0 + plot_w, nominal_y), fill=hex_color("gold"), width=3)
        draw.text((x0 + 230, nominal_y - 25), "nominal 0.95", fill=hex_color("gold"), font=tiny)
        for tick in [0.5, 0.7, 0.9, 1.0]:
            y = y0 - int((tick - 0.45) / 0.55 * plot_h)
            draw.line((x0 - 4, y, x0, y), fill=hex_color("muted"), width=1)
            draw.text((x0 - 48, y - 9), f"{tick:.1f}", fill=hex_color("muted"), font=tiny)
        for i, rho in enumerate(ICC_GRID):
            x = x0 + 44 + i * 90
            draw.line((x, y0, x, y0 + 5), fill=hex_color("muted"), width=1)
            draw.text((x - 12, y0 + 14), f"{rho:.1f}", fill=hex_color("muted"), font=tiny)
        for method in ["naive_iid_ols_z", "cluster_robust_z"]:
            points = []
            for i, rho in enumerate(ICC_GRID):
                row = next(item for item in rows if item["G"] == g_count and item["rho"] == rho and item["imbalance"] == "imbalanced")
                value = row["methods"][method]["coverage"]
                mc = row["methods"][method]["mc_se"]
                x = x0 + 44 + i * 90
                y = y0 - int((value - 0.45) / 0.55 * plot_h)
                lo = y0 - int((value - 1.96 * mc - 0.45) / 0.55 * plot_h)
                hi = y0 - int((value + 1.96 * mc - 0.45) / 0.55 * plot_h)
                draw.line((x, hi, x, lo), fill=method_colors[method], width=2)
                draw.ellipse((x - 7, y - 7, x + 7, y + 7), fill=method_colors[method])
                points.append((x, y))
            draw.line(points, fill=method_colors[method], width=4)
    stress = summary["negative_result"]
    x_arrow = 92 + 44 + 3 * 90
    y_arrow = 575 - int((stress["cluster_robust_coverage"] - 0.45) / 0.55 * 410)
    draw.line((x_arrow + 52, y_arrow - 88, x_arrow + 10, y_arrow - 12), fill=hex_color("purple"), width=3)
    draw.polygon([(x_arrow + 10, y_arrow - 12), (x_arrow + 22, y_arrow - 16), (x_arrow + 15, y_arrow - 26)], fill=hex_color("purple"))
    draw.text((x_arrow + 58, y_arrow - 112), "small-G stress", fill=hex_color("purple"), font=small)
    draw.text((x_arrow + 58, y_arrow - 88), f"cluster={stress['cluster_robust_coverage']:.2f}", fill=hex_color("purple"), font=tiny)
    draw.rectangle((1115, 28, 1142, 49), fill=method_colors["naive_iid_ols_z"])
    draw.text((1152, 25), "naive iid OLS z", fill=hex_color("ink"), font=small)
    draw.rectangle((1115, 60, 1142, 81), fill=method_colors["cluster_robust_z"])
    draw.text((1152, 57), "cluster-robust z", fill=hex_color("ink"), font=small)
    img.save(path)


def draw_negative_plot(path: Path, summary: dict[str, object]) -> None:
    img = Image.new("RGB", (1180, 520), "white")
    draw = ImageDraw.Draw(img)
    font = load_font(22)
    small = load_font(16)
    title = load_font(27)
    draw.text((46, 28), "G=8 with imbalanced clusters: coverage by ICC", fill=hex_color("ink"), font=title)
    x0, y0 = 112, 425
    draw.line((x0, 88, x0, y0), fill=hex_color("muted"), width=2)
    draw.line((x0, y0, 1060, y0), fill=hex_color("muted"), width=2)
    nominal_y = y0 - int((0.95 - 0.45) / 0.55 * 310)
    draw.line((x0, nominal_y, 1060, nominal_y), fill=hex_color("gold"), width=3)
    draw.text((900, nominal_y - 28), "target 0.95", fill=hex_color("gold"), font=small)
    rows = summary["rows"]
    colors = {"naive_iid_ols_z": hex_color("red"), "cluster_robust_z": hex_color("teal")}
    for i, rho in enumerate(ICC_GRID):
        row = next(item for item in rows if item["G"] == 8 and item["rho"] == rho and item["imbalance"] == "imbalanced")
        cx = x0 + 105 + i * 210
        for j, method in enumerate(["naive_iid_ols_z", "cluster_robust_z"]):
            value = row["methods"][method]["coverage"]
            bar_h = int((value - 0.45) / 0.55 * 310)
            x = cx + j * 50
            draw.rectangle((x, y0 - bar_h, x + 36, y0), fill=colors[method])
            draw.text((x - 2, y0 - bar_h - 28), f"{value:.2f}", fill=hex_color("ink"), font=small)
        draw.text((cx + 4, y0 + 20), f"ρ={rho:.1f}", fill=hex_color("muted"), font=font)
    final_x = x0 + 105 + 3 * 210 + 50
    stress_y = y0 - int((summary["negative_result"]["cluster_robust_coverage"] - 0.45) / 0.55 * 310)
    draw.line((final_x + 105, stress_y - 92, final_x + 43, stress_y - 14), fill=hex_color("purple"), width=3)
    draw.polygon([(final_x + 43, stress_y - 14), (final_x + 56, stress_y - 18), (final_x + 49, stress_y - 30)], fill=hex_color("purple"))
    draw.text((final_x + 112, stress_y - 114), "still below nominal", fill=hex_color("purple"), font=small)
    draw.rectangle((790, 32, 818, 53), fill=colors["naive_iid_ols_z"])
    draw.text((828, 29), "naive", fill=hex_color("ink"), font=small)
    draw.rectangle((790, 65, 818, 86), fill=colors["cluster_robust_z"])
    draw.text((828, 62), "cluster-robust", fill=hex_color("ink"), font=small)
    img.save(path)


def draw_model_page(slide, assets: Path, manifest: dict, refs: list[str], summary: dict[str, object]):
    audience_text = [
        "Synthetic multi-center simulation; treatment effect β₁=0.25.",
        "β₁ is the treatment effect.",
        "uⱼ: center-level shift shared within center j.",
        "εᵢⱼ: individual residual noise.",
        "ρ is the intraclass correlation; center j is the inference unit.",
        "Point estimates can remain centered while iid intervals lose coverage.",
    ]
    manifest.setdefault("audience_text_by_archetype", {})["STATISTICAL_MODEL"] = audience_text
    add_text(slide, audience_text[0], 0.82, 1.16, 11.8, 0.26, 11.2, "muted", True)
    add_math(
        slide,
        assets,
        manifest,
        "slide1_dgp",
        r"$Y_{ij}=\beta_0+\beta_1T_{ij}+u_j+\varepsilon_{ij}$",
        0.85,
        1.74,
        7.15,
        fontsize=42,
    )
    add_math(
        slide,
        assets,
        manifest,
        "slide1_components",
        r"$u_j\sim N(0,\tau^2),\qquad \varepsilon_{ij}\sim N(0,\sigma^2)$",
        0.92,
        3.02,
        6.85,
        fontsize=31,
    )
    add_math(
        slide,
        assets,
        manifest,
        "slide1_icc",
        r"$\rho=\frac{\tau^2}{\tau^2+\sigma^2}$",
        8.72,
        2.0,
        2.45,
        fontsize=44,
        color="purple",
    )
    add_text(slide, audience_text[1], 0.98, 4.12, 2.6, 0.28, 13.2, "blue", True)
    add_text(slide, audience_text[2], 3.7, 4.12, 3.35, 0.35, 12.2, "teal", True)
    add_text(slide, audience_text[3], 7.15, 4.12, 2.65, 0.35, 12.2, "gold", True)
    add_text(slide, audience_text[4], 8.55, 3.05, 3.2, 0.62, 12.4, "purple", True)
    add_rule(slide, 0.95, 4.82, 11.75, 4.82, "line", 1.0)
    add_text(slide, audience_text[5], 1.0, 5.14, 8.9, 0.48, 15.2, "ink", True)
    add_text(slide, "Grid: G=8/20/50 centers, ICC ρ=0/.1/.3/.5, balanced and imbalanced cluster sizes.", 1.0, 5.78, 10.9, 0.32, 11.6, "muted")
    manifest["model_page_structure"] = {
        "additive_components": ["u_j", "epsilon_ij"],
        "connector_policy": "center random effect and individual residual are parallel additive components; notation is grounded by proximity and labels",
        "math_asset_keys": ["slide1_dgp", "slide1_components", "slide1_icc"],
    }


def draw_estimator_page(slide, assets: Path, manifest: dict, refs: list[str], summary: dict[str, object]):
    audience_text = [
        "Both intervals estimate β₁; only the variance object changes.",
        "Cluster meat: score contributions are summed by center before forming covariance.",
        "Naive variance treats rows as independent.",
        "The simulation uses cluster-robust z intervals; CR2 and wild bootstrap remain planned comparisons.",
    ]
    manifest.setdefault("audience_text_by_archetype", {})["ESTIMATOR"] = audience_text
    add_text(slide, audience_text[0], 0.85, 1.16, 11.5, 0.28, 11.8, "muted", True)
    add_math(
        slide,
        assets,
        manifest,
        "slide2_sandwich",
        r"$\widehat V_{\mathrm{CR}}=(X^\top X)^{-1}\left(\sum_g X_g^\top \hat u_g\hat u_g^\top X_g\right)(X^\top X)^{-1}$",
        0.76,
        1.78,
        10.1,
        fontsize=36,
        color="ink",
    )
    add_rule(slide, 5.15, 3.0, 7.48, 3.0, "purple", 3.0)
    arrow(slide, 6.45, 3.08, 6.45, 3.72, "purple", 1.8)
    add_text(slide, audience_text[1], 4.15, 3.86, 4.9, 0.46, 13.0, "purple", True, PP_ALIGN.CENTER)
    add_math(
        slide,
        assets,
        manifest,
        "slide2_naive",
        r"$\widehat V_{\mathrm{iid}}=\hat\sigma^2(X^\top X)^{-1}$",
        0.98,
        4.98,
        3.2,
        fontsize=30,
        color="red",
    )
    add_text(slide, audience_text[2], 1.05, 5.76, 3.7, 0.26, 11.6, "red", True)
    add_text(slide, audience_text[3], 5.28, 5.03, 6.55, 0.58, 13.4, "ink", True)
    manifest["estimator_page_structure"] = {
        "math_asset_keys": ["slide2_sandwich", "slide2_naive"],
        "annotated_term": "sum_g X_g^T u_g u_g^T X_g",
        "source_like_formula_text_in_audience_objects": False,
    }


def draw_design_page(slide, assets: Path, manifest: dict, refs: list[str], summary: dict[str, object]):
    audience_text = [
        "Synthetic experiment: vary center count, ICC, and cluster-size balance.",
        f"{REPLICATES} replicates per cell",
        "naive iid OLS z interval",
        "cluster-robust z interval",
        "coverage, bias, interval width",
        "The design asks whether center aggregation repairs interval coverage across the stress grid.",
    ]
    manifest.setdefault("audience_text_by_archetype", {})["SIMULATION_DESIGN"] = audience_text
    add_text(slide, audience_text[0], 0.85, 1.16, 11.5, 0.28, 11.8, "muted", True)
    rect(slide, "DGP knobs\nG=8,20,50 centers\nICC ρ=0,.1,.3,.5\nbalanced vs imbalanced", 0.82, 2.08, 2.55, 1.36, "white", "ink", 11.8, True, line="gold")
    rect(slide, f"Generated samples\n{audience_text[1]}\noutcome, treatment, center", 3.94, 2.08, 2.55, 1.36, "white", "ink", 11.8, True, line="blue")
    rect(slide, audience_text[2], 7.05, 1.64, 2.36, 0.58, "white", "red", 11.8, True, line="red")
    rect(slide, audience_text[3] + "\ncenter scores aggregated", 7.05, 3.05, 2.36, 0.78, "white", "teal", 11.0, True, line="teal")
    rect(slide, audience_text[4], 10.08, 2.33, 2.3, 0.94, "white", "ink", 11.8, True, line="purple")
    arrow(slide, 3.38, 2.76, 3.9, 2.76, "muted", 1.6)
    arrow(slide, 6.5, 2.76, 7.02, 1.94, "muted", 1.6)
    arrow(slide, 6.5, 2.76, 7.02, 3.42, "muted", 1.6)
    arrow(slide, 9.43, 1.94, 10.05, 2.62, "muted", 1.6)
    arrow(slide, 9.43, 3.42, 10.05, 2.98, "muted", 1.6)
    add_rule(slide, 0.95, 4.82, 11.85, 4.82, "line", 1.0)
    add_text(slide, audience_text[5], 1.0, 5.18, 10.8, 0.48, 14.0, "ink", True)
    add_text(slide, "Simulation study; endpoints are computed from the deterministic output for the same DGP grid.", 1.0, 5.84, 10.8, 0.32, 11.5, "muted")
    manifest["simulation_design_diagram"] = {
        "structural_connectors": ["dgp_to_data", "data_to_naive", "data_to_cluster", "naive_to_endpoint", "cluster_to_endpoint"],
        "arrowheads": "DrawingML a:tailEnd type=triangle on every connector",
        "edge_crossing": "none",
        "reading_direction": "left_to_right",
        "endpoint_semantics": ["95% interval coverage", "bias", "interval width"],
    }


def draw_result_page(slide, assets: Path, manifest: dict, refs: list[str], summary: dict[str, object]):
    audience_text = [
        "Synthetic imbalanced-cluster simulation; bands show Monte Carlo uncertainty.",
        "Cluster-robust intervals recover much of the lost coverage, but small-G stress remains below nominal.",
    ]
    manifest.setdefault("audience_text_by_archetype", {})["RESULT_FIGURE"] = audience_text
    plot = assets / "coverage_by_icc.png"
    draw_coverage_plot(plot, summary)
    slide.shapes.add_picture(str(plot), inch(0.82), inch(1.28), width=inch(10.9))
    add_text(slide, audience_text[0], 0.95, 6.42, 8.9, 0.28, 10.5, "muted")


def draw_negative_page(slide, assets: Path, manifest: dict, refs: list[str], summary: dict[str, object]):
    stress = summary["negative_result"]
    audience_text = [
        f"At G=8, ICC ρ=.5, imbalanced clusters, cluster-robust coverage is {stress['cluster_robust_coverage']:.2f}.",
        "Few centers make the sandwich variance noisy; imbalance increases center leverage.",
        "Planned comparison: CR2 small-sample correction and wild cluster bootstrap on the same grid.",
    ]
    manifest.setdefault("audience_text_by_archetype", {})["NEGATIVE_RESULT"] = audience_text
    plot = assets / "small_g_negative_result.png"
    draw_negative_plot(plot, summary)
    slide.shapes.add_picture(str(plot), inch(0.58), inch(1.38), width=inch(8.3))
    add_text(slide, audience_text[0], 9.25, 1.55, 2.75, 0.8, 13.0, "red", True)
    add_rule(slide, 9.25, 2.58, 11.85, 2.58, "line", 1.0)
    add_text(slide, audience_text[1], 9.25, 2.86, 2.9, 0.72, 12.1, "ink", True)
    add_rule(slide, 9.25, 3.82, 11.85, 3.82, "line", 1.0)
    add_text(slide, audience_text[2], 9.25, 4.10, 2.9, 0.86, 11.8, "teal", True)
    add_text(slide, "Completed evidence is limited to the two intervals shown; the correction methods are future tests.", 0.95, 6.2, 11.1, 0.32, 11.8, "muted")
    manifest["negative_result_claim"] = stress


SLIDES = [
    ("Center correlation breaks iid intervals", "The point estimate targets β₁, but the uncertainty model changes.", "STATISTICAL_MODEL", draw_model_page),
    ("The sandwich variance aggregates by center", "The middle term is a center-level score covariance, not a row-level residual pool.", "ESTIMATOR", draw_estimator_page),
    ("A stress grid separates interval behavior", "The same generated samples feed two intervals and the same coverage endpoint.", "SIMULATION_DESIGN", draw_design_page),
    ("Cluster adjustment recovers coverage except at small G", "Coverage near 0.95 is the target; uncertainty is part of the figure.", "RESULT_FIGURE", draw_result_page),
    ("The remaining failure is finite-center undercoverage", "The next experiment tests small-sample corrections under the same grid.", "NEGATIVE_RESULT", draw_negative_page),
]


INTERNAL_LEAK_PATTERNS = [
    "RRL-",
    "Reference retrieval",
    "EVIDENCE_MANIFEST",
    "Diagram contract",
    "style not copied",
    "Reading target",
    "Observed in this synthetic run",
    "repo path",
    "run ID",
    "implementation commit",
    "review target",
]

ANTI_META_PHRASES = [
    "Evidence boundary",
    "Key change",
    "Role in the deck",
    "This slide",
    "This page",
]

MATH_SOURCE_PATTERNS = [
    "beta_",
    "epsilon_",
    "sum_",
    "^(-1)",
    "X'X",
    "rho =",
]


def build_pptx(path: Path, assets: Path, manifest: dict, summary: dict[str, object]) -> None:
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)
    blank = prs.slide_layouts[6]
    manifest["slides"] = []
    manifest["reference_design_audit"] = {}
    for index, (title, message, archetype, drawer) in enumerate(SLIDES, start=1):
        retrieval = retrieve_references(archetype)
        refs = list(retrieval["selected_ids"])
        slide = prs.slides.add_slide(blank)
        header(slide, index, title, message)
        drawer(slide, assets, manifest, refs, summary)
        audit = reference_design_audit(archetype, retrieval)
        manifest["reference_design_audit"][f"slide_{index}"] = audit
        manifest["slides"].append({
            "slide": index,
            "title": title,
            "archetype": archetype,
            "reference_ids": refs,
            "reference_retrieval": retrieval,
            "reference_design_audit": audit,
            "audience_text": manifest.get("audience_text_by_archetype", {}).get(archetype, []),
            "learned_organization": retrieval["organization_lesson"],
            "reference_rationale": "References are selected by auditable semantic query over inspected page records with rendered-page checksums; they inform organization only.",
            "style_not_copied": "No whole-slide screenshot, public slide styling, private CARE figure, or clinical/patient image is copied.",
            "what_not_copied": retrieval["what_was_not_copied"],
            "evidence_status": "synthetic_simulation_generated" if archetype in {"RESULT_FIGURE", "NEGATIVE_RESULT"} else "synthetic_model_or_simulation_context",
            "expected_scientific_objects": {
                "STATISTICAL_MODEL": ["beta_1 estimand", "DGP formula", "u_j random effect", "epsilon_ij residual", "ICC rho", "center inference unit"],
                "ESTIMATOR": ["naive variance formula", "cluster sandwich formula", "X_g block", "u_g residual vector", "center aggregation"],
                "SIMULATION_DESIGN": ["DGP knobs", "replicates", "naive interval branch", "cluster-robust interval branch", "coverage endpoint gate"],
                "RESULT_FIGURE": ["coverage plot", "nominal 95% line", "method comparison", "Monte Carlo uncertainty", "synthetic boundary"],
                "NEGATIVE_RESULT": ["small-G stress evidence", "coverage shortfall", "failure mechanism", "planned CR2/wild bootstrap experiment"],
            }[archetype],
        })
    prs.save(path)


def pptx_slide_text(pptx_path: Path) -> str:
    with ZipFile(pptx_path) as zf:
        slide_names = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        return "\n".join(zf.read(name).decode("utf-8", errors="ignore") for name in slide_names)


def run_quality_gates(pptx_path: Path, manifest: dict[str, object]) -> dict[str, object]:
    slide_xml = pptx_slide_text(pptx_path)
    audience_text = "\n".join(
        "\n".join(slide.get("audience_text", []))
        for slide in manifest.get("slides", [])
        if isinstance(slide, dict)
    )
    failures: list[str] = []
    for pattern in INTERNAL_LEAK_PATTERNS:
        if pattern in slide_xml or pattern in audience_text:
            failures.append(f"audience-facing internal leak: {pattern}")
    for phrase in ANTI_META_PHRASES:
        if phrase in slide_xml or phrase in audience_text:
            failures.append(f"audience-facing meta language: {phrase}")
    math_pages = "\n".join(
        "\n".join(slide.get("audience_text", []))
        for slide in manifest.get("slides", [])
        if isinstance(slide, dict) and slide.get("archetype") in {"STATISTICAL_MODEL", "ESTIMATOR"}
    )
    for pattern in MATH_SOURCE_PATTERNS:
        if pattern in math_pages or pattern in slide_xml:
            failures.append(f"source-like math leaked into audience text: {pattern}")
    math_assets = manifest.get("math_assets", {})
    for required in ["slide1_dgp", "slide1_components", "slide1_icc", "slide2_sandwich", "slide2_naive"]:
        asset = math_assets.get(required) if isinstance(math_assets, dict) else None
        if not isinstance(asset, dict) or not Path(str(asset.get("path", ""))).is_file():
            failures.append(f"missing rendered math asset: {required}")
    if "u_j_to_epsilon_ij" in json.dumps(manifest.get("model_page_structure", {})):
        failures.append("model page still records serial u_j_to_epsilon_ij connector")
    audit = manifest.get("reference_design_audit", {})
    if not isinstance(audit, dict) or len(audit) != 5:
        failures.append("reference_design_audit must cover five slides")
    audit_items = audit.items() if isinstance(audit, dict) else []
    for slide_key, item in audit_items:
        if len(item.get("selected_reference_ids", [])) < 2:
            failures.append(f"{slide_key} has fewer than two inspected reference IDs")
        if not item.get("adopted_design_decisions"):
            failures.append(f"{slide_key} lacks adopted design decisions")
    return {
        "status": "PASS" if not failures else "FAIL",
        "failures": failures,
        "checked_gates": [
            "audience_facing_internal_leak",
            "math_source_leak",
            "rendered_math_assets_present",
            "anti_meta_language",
            "reference_design_audit",
        ],
    }


def editable_slide_count(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])


def find_renderer() -> str | None:
    explicit = os.environ.get("RESEARCH_PPTX_RENDERER")
    if explicit and Path(explicit).exists():
        return explicit
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    tools = Path(".cache/tools")
    if tools.exists():
        for extracted_soffice in sorted(tools.glob("squashfs-root/opt/libreoffice*/program/soffice")):
            return str(extracted_soffice)
        extracted = tools / "squashfs-root" / "AppRun"
        if extracted.exists():
            return str(extracted)
        for candidate in sorted(tools.glob("LibreOffice*.AppImage")):
            return str(candidate)
    return None


def render_pptx(pptx_path: Path, out_dir: Path) -> dict:
    renderer = find_renderer()
    render_dir = out_dir / "rendered"
    pdf_dir = out_dir / "pdf"
    profile_dir = out_dir / "lo-profile"
    home_dir = out_dir / "lo-home"
    cache_dir = out_dir / "fontconfig-cache"
    render_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    home_dir.mkdir(parents=True, exist_ok=True)
    cache_dir.mkdir(parents=True, exist_ok=True)
    if not renderer:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "reason": "no soffice/libreoffice/RESEARCH_PPTX_RENDERER found", "png_count": 0}
    profile_uri = profile_dir.resolve().as_uri()
    env = os.environ.copy()
    env["HOME"] = str(home_dir.resolve())
    env["XDG_CACHE_HOME"] = str(cache_dir.resolve())
    cmd = [
        renderer,
        f"-env:UserInstallation={profile_uri}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir.resolve()),
        str(pptx_path.resolve()),
    ]
    result = subprocess.run(cmd, check=False, capture_output=True, text=True, timeout=120, env=env)
    pdf_candidates = sorted(pdf_dir.glob("*.pdf"))
    if result.returncode != 0 or not pdf_candidates:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "returncode": result.returncode, "stderr": result.stderr[-1000:], "png_count": 0}
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pdftoppm not found after real PPTX-to-PDF conversion", "png_count": 0}
    prefix = render_dir / "slide"
    ppm = subprocess.run([pdftoppm, "-png", "-r", "160", str(pdf_candidates[0]), str(prefix)], check=False, capture_output=True, text=True, timeout=120)
    pngs = sorted(render_dir.glob("slide-*.png"))
    return {
        "status": "ok" if ppm.returncode == 0 and len(pngs) == 5 else "BLOCKED_REAL_PPTX_RENDER",
        "renderer": renderer,
        "pptx_to_pdf": str(pdf_candidates[0]),
        "png_count": len(pngs),
        "rendered_pngs": [str(path) for path in pngs],
        "returncode": ppm.returncode,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out-dir", type=Path, default=Path(".cache/statistical-method-group-meeting-benchmark"))
    args = parser.parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)
    assets = args.out_dir / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    summary = run_simulation(args.out_dir)
    pptx_path = args.out_dir / "statistical_method_group_meeting_benchmark.pptx"
    manifest_path = args.out_dir / "EVIDENCE_MANIFEST.json"
    reference_audit_path = args.out_dir / "reference_design_audit.json"
    render_status_path = args.out_dir / "RENDER_STATUS.json"
    manifest: dict = {
        "status": "GENERATED_SOURCE_ARTIFACTS_ONLY",
        "generator_may_pass": False,
        "task_key": "016_statistical_method_group_meeting_benchmark",
        "rights_note": "Synthetic data only; no downloaded public deck image, private CARE figure, whole-slide screenshot, or patient image is copied.",
        "benchmark_story": "Multi-center correlation can leave beta_1 estimates nearly unbiased while interval coverage fails; cluster-robust z intervals repair much of the issue but small-G/high-ICC stress remains a negative result.",
        "simulation_summary": summary,
    }
    build_pptx(pptx_path, assets, manifest, summary)
    manifest["pptx"] = str(pptx_path)
    manifest["editable_slide_count"] = editable_slide_count(pptx_path)
    manifest["deterministic_quality_gates"] = run_quality_gates(pptx_path, manifest)
    if manifest["deterministic_quality_gates"]["status"] != "PASS":
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        raise SystemExit("deterministic presentation QA failed: " + "; ".join(manifest["deterministic_quality_gates"]["failures"]))
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reference_audit_path.write_text(json.dumps(manifest["reference_design_audit"], ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    render = render_pptx(pptx_path, args.out_dir)
    render_status_path.write_text(json.dumps(render, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"pptx": str(pptx_path), "evidence_manifest": str(manifest_path), "render_status": str(render_status_path)}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
