#!/usr/bin/env python3
"""Generate a four-slide research-group-meeting regression deck.

The generator creates only source artifacts and evidence. It never writes a
final scientific PASS; that decision belongs to the independent reviewer.
"""

from __future__ import annotations

import argparse
import csv
import json
import math
import os
import shutil
import subprocess
from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


W, H = 13.333, 7.5
REPO_ROOT = Path(__file__).resolve().parents[4]
REFERENCE_ROOT = REPO_ROOT / "skills" / "tools" / "documents-media" / "presentations" / "shared" / "references"
REFERENCE_INDEX = REFERENCE_ROOT / "research_slide_reference_index.csv"
REFERENCE_MANIFEST = REFERENCE_ROOT / "reference_sources_manifest.json"
P = {
    "bg": "F7F7F9",
    "ink": "17202A",
    "muted": "606977",
    "line": "C6CCD6",
    "purple": "4F1F68",
    "teal": "0F766E",
    "blue": "1F4E79",
    "gold": "9A6A16",
    "red": "A33A34",
    "soft_teal": "E1F4F1",
    "soft_blue": "E5EEF8",
    "soft_gold": "FBF1D6",
    "soft_red": "F8E3E1",
}

REFERENCE_QUERIES = {
    "RESULT_FIGURE": {
        "intent": "Show endpoint-wise quantitative results where interpretation changes with metric choice.",
        "page_functions": ["RESULT_FIGURE", "CONFIDENCE_INTERVAL", "SENSITIVITY_ANALYSIS", "REAL_DATA_APPLICATION"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "annotation variability", "survey nonresponse / MRP", "Bayesian workflow"],
        "evidence_types": ["quantitative plot", "quantitative table", "uncertainty interval", "time-series interval", "subgroup comparison"],
        "organization_lesson": "Put the endpoint, comparator, uncertainty, and decision implication adjacent to the main result graphic.",
    },
    "FAILURE_CASE": {
        "intent": "Show a single synthetic hard case with aligned image, GT, prediction, error overlay, and metrics.",
        "page_functions": ["MEDICAL_IMAGE_COMPARISON", "NEGATIVE_RESULT", "MODEL_CHECK", "RESULT_FIGURE"],
        "scientific_domain": ["medical imaging", "statistics"],
        "statistical_subdomain": ["lesion segmentation", "annotation variability", "Bayesian workflow"],
        "evidence_types": ["image comparison", "model comparison", "negative/fix title", "posterior predictive check"],
        "organization_lesson": "Keep the failing object visible next to the metric so the audience can inspect the mechanism, not only the score.",
    },
    "EXPERIMENT_DESIGN": {
        "intent": "Explain a multi-center experiment with local data, transmitted summaries, estimator, comparator, and endpoint gate.",
        "page_functions": ["EXPERIMENT_DESIGN", "METHOD_DIAGRAM", "ESTIMATOR", "NEXT_EXPERIMENT", "ASSUMPTION"],
        "scientific_domain": ["medical imaging", "statistics", "research communication"],
        "statistical_subdomain": ["lesion segmentation", "hybrid resource-bound analysis", "doctoral update", "survey nonresponse / MRP"],
        "evidence_types": ["task overview", "method mechanism", "estimator pipeline", "planned evidence", "assumption / challenge"],
        "organization_lesson": "Connect units, local estimators, shared summaries, and evaluation endpoints in one inspectable flow.",
    },
    "STATISTICAL_MODEL": {
        "intent": "Show a measured-variable inference mechanism where AI phenotype error changes the target estimand.",
        "page_functions": ["STATISTICAL_MODEL", "BAYESIAN_MODEL", "ESTIMATOR", "MODEL_CHECK", "POSTERIOR_DIAGNOSTIC"],
        "scientific_domain": ["statistics", "biostatistics", "medical imaging"],
        "statistical_subdomain": ["lesion segmentation", "Bayesian workflow", "Bayesian priors", "Bayesian data analysis", "survey nonresponse / MRP"],
        "evidence_types": ["model objective", "Bayesian model", "estimator formula", "modeling workflow", "posterior diagnostic"],
        "organization_lesson": "Name the estimand, observed proxy, validation subset, model correction, and inference boundary explicitly.",
    },
}

TIER_PRIORITY = {
    "PRIMARY_RESEARCH_PRESENTATION": 0,
    "SECONDARY_TEACHING_REFERENCE": 1,
    "PRESENTATION_GUIDANCE": 2,
    "CANDIDATE_BACKLOG": 3,
}


def rgb(name: str) -> RGBColor:
    value = P[name]
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def hex_color(name: str) -> str:
    return "#" + P[name]


def inch(value: float):
    return Inches(value)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float = 13, color: str = "ink", bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def rect(slide, text: str, x: float, y: float, w: float, h: float, fill: str = "soft_blue", color: str = "ink", size: float = 12, bold: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("line")
    shape.line.width = Pt(1)
    add_text(slide, text, x + 0.06, y + 0.05, w - 0.12, h - 0.1, size, color, bold, PP_ALIGN.CENTER)
    return shape


def line(slide, x1: float, y1: float, x2: float, y2: float, color: str = "muted", width: float = 1.4):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    c.line.color.rgb = rgb(color)
    c.line.width = Pt(width)
    return c


def swatch(slide, x: float, y: float, color: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(0.12), inch(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(0.5)
    return shape


def header(slide, number: int, title: str, message: str):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0), inch(0), inch(W), inch(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb("purple")
    band.line.fill.background()
    add_text(slide, title, 0.55, 0.38, 9.5, 0.36, 21, "ink", True)
    add_text(slide, f"{number:02d}/04", 11.6, 0.42, 0.95, 0.22, 9, "muted", False, PP_ALIGN.RIGHT)
    add_text(slide, message, 0.65, 0.88, 11.7, 0.34, 12.4, "purple", True)


def tokens(value: str) -> set[str]:
    return {part.lower() for part in value.replace("/", " ").replace("-", " ").replace("_", " ").split() if len(part) > 2}


def load_reference_rows() -> list[dict[str, str]]:
    manifest = json.loads(REFERENCE_MANIFEST.read_text(encoding="utf-8"))
    source_tiers = {item["source_id"]: item["source_tier"] for item in manifest["candidate_sources"]}
    source_domains = {item["source_id"]: item["domain_family"] for item in manifest["candidate_sources"]}
    source_subdomains = {item["source_id"]: item["statistical_subdomain"] for item in manifest["candidate_sources"]}
    rows = []
    with REFERENCE_INDEX.open(encoding="utf-8", newline="") as fh:
        for row in csv.DictReader(fh):
            if row.get("verification_status") != "inspected":
                continue
            if not (row.get("source_file_sha256") and row.get("rendered_page_sha256")):
                continue
            if not (row.get("inspection_date") and row.get("inspection_means")):
                continue
            row["source_tier"] = source_tiers.get(row["source_id"], "")
            row["domain_family"] = source_domains.get(row["source_id"], "")
            row["statistical_subdomain"] = source_subdomains.get(row["source_id"], "")
            rows.append(row)
    return rows


def score_reference(row: dict[str, str], query: dict[str, object]) -> tuple[int, list[str]]:
    score = 0
    reasons: list[str] = []
    page_functions = query["page_functions"]
    evidence_types = query["evidence_types"]
    domains = query["scientific_domain"]
    subdomains = query["statistical_subdomain"]
    if row["page_function"] in page_functions:
        score += 10
        reasons.append(f"page_function={row['page_function']} matches query")
    if row["evidence_type"] in evidence_types:
        score += 6
        reasons.append(f"evidence_type={row['evidence_type']} matches query")
    if row["domain_family"] in domains:
        score += 4
        reasons.append(f"domain={row['domain_family']} matches query")
    if row["statistical_subdomain"] in subdomains:
        score += 4
        reasons.append(f"subdomain={row['statistical_subdomain']} matches query")
    query_terms = tokens(" ".join([str(query["intent"]), " ".join(page_functions), " ".join(evidence_types)]))
    row_terms = tokens(" ".join([
        row["scientific_object"],
        row["why_this_specific_page_works"],
        row["what_to_learn"],
        row["short_page_specific_observation"],
    ]))
    overlap = sorted(query_terms & row_terms)
    if overlap:
        score += min(5, len(overlap))
        reasons.append("semantic tokens overlap: " + ", ".join(overlap[:6]))
    tier_bonus = max(0, 3 - TIER_PRIORITY.get(row["source_tier"], 3))
    score += tier_bonus
    reasons.append(f"source_tier={row['source_tier']} priority_bonus={tier_bonus}")
    return score, reasons


def retrieve_references(archetype: str) -> dict[str, object]:
    query = REFERENCE_QUERIES[archetype]
    scored = []
    for row in load_reference_rows():
        score, reasons = score_reference(row, query)
        if score <= 0:
            continue
        scored.append((score, TIER_PRIORITY.get(row["source_tier"], 3), row["reference_id"], row, reasons))
    scored.sort(key=lambda item: (-item[0], item[1], item[2]))
    candidate_rows = scored[:8]
    selected_rows = candidate_rows[:5]
    if len(selected_rows) < 2:
        raise RuntimeError(f"reference retrieval for {archetype} found fewer than two inspected references")
    return {
        "query": query,
        "candidate_ids": [row["reference_id"] for _, _, _, row, _ in candidate_rows],
        "selected_ids": [row["reference_id"] for _, _, _, row, _ in selected_rows],
        "source_tiers": {row["reference_id"]: row["source_tier"] for _, _, _, row, _ in selected_rows},
        "ranking_relevance_reason": {
            row["reference_id"]: "; ".join(reasons)
            for _, _, _, row, reasons in selected_rows
        },
        "organization_lesson": query["organization_lesson"],
        "what_was_not_copied": "No full-slide screenshots, source images, institutional styling, public figures, private clinical data, or source-specific visual identity were copied.",
    }


def reference_footer(refs: list[str], purpose: str) -> str:
    return f"Reference retrieval: {purpose} selected inspected pages {', '.join(refs)}; trace in EVIDENCE_MANIFEST; style not copied."


ENDPOINT_DATA = {
    "Dice": {"Baseline": (0.82, 0.03), "Calibrated": (0.79, 0.04), "Federated": (0.76, 0.05), "Local-only": (0.72, 0.06)},
    "Lesion recall": {"Baseline": (0.58, 0.07), "Calibrated": (0.74, 0.05), "Federated": (0.69, 0.06), "Local-only": (0.63, 0.07)},
    "Burden error": {"Baseline": (0.31, 0.05), "Calibrated": (0.18, 0.04), "Federated": (0.22, 0.05), "Local-only": (0.35, 0.06)},
}


def load_font(size: int):
    for path in ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/dejavu/DejaVuSans.ttf"]:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_endpoint_chart(path: Path) -> dict:
    img = Image.new("RGB", (1060, 520), "white")
    draw = ImageDraw.Draw(img)
    font = load_font(18)
    small = load_font(14)
    colors = {"Baseline": "#1F4E79", "Calibrated": "#0F766E", "Federated": "#9A6A16", "Local-only": "#A33A34"}
    x0, y0 = 90, 390
    draw.line((x0, 55, x0, y0), fill="#5F6772", width=2)
    draw.line((x0, y0, 1010, y0), fill="#5F6772", width=2)
    for tick in range(0, 101, 20):
        y = y0 - tick * 3
        draw.line((x0 - 5, y, x0, y), fill="#5F6772", width=1)
        draw.text((38, y - 8), f"{tick/100:.1f}", fill="#17202A", font=small)
        draw.line((x0, y, 1010, y), fill="#E2E6EC", width=1)
    endpoints = list(ENDPOINT_DATA)
    methods = list(next(iter(ENDPOINT_DATA.values())))
    bar_w = 34
    group_w = 270
    for i, endpoint in enumerate(endpoints):
        gx = x0 + 60 + i * group_w
        endpoint_label = endpoint if endpoint != "Burden error" else "Burden error\nlower is better"
        draw.multiline_text((gx + 30, 420), endpoint_label, fill="#17202A", font=font, spacing=2)
        for j, method in enumerate(methods):
            value, err = ENDPOINT_DATA[endpoint][method]
            score = value
            x = gx + j * (bar_w + 12)
            y = y0 - int(score * 300)
            draw.rectangle((x, y, x + bar_w, y0), fill=colors[method])
            err_px = int(err * 300)
            draw.line((x + bar_w // 2, y - err_px, x + bar_w // 2, y + err_px), fill="#17202A", width=2)
            draw.text((x - 3, y - 25), f"{value:.2f}", fill="#17202A", font=small)
    for j, method in enumerate(methods):
        lx = 690 + (j % 2) * 170
        ly = 35 + (j // 2) * 24
        draw.rectangle((lx, ly, lx + 18, ly + 18), fill=colors[method])
        draw.text((lx + 25, ly - 1), method, fill="#17202A", font=small)
    img.save(path)
    best_by_endpoint = {
        endpoint: min(values, key=lambda key: values[key][0]) if endpoint == "Burden error" else max(values, key=lambda key: values[key][0])
        for endpoint, values in ENDPOINT_DATA.items()
    }
    return {
        "data": ENDPOINT_DATA,
        "best_by_endpoint": best_by_endpoint,
        "burden_error_favorable_direction": "lower_is_better",
        "display_encoding": {"Burden error": "raw_error_value_lower_is_better"},
    }


def inside_ellipse(x: int, y: int, cx: int, cy: int, rx: int, ry: int) -> bool:
    return ((x - cx) / rx) ** 2 + ((y - cy) / ry) ** 2 <= 1


def draw_phantom(path: Path) -> dict:
    w, h = 1040, 430
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    font = load_font(17)
    small = load_font(13)
    panel_w = 235
    source_size = 120
    rendered_case_size = 188
    labels = ["Synthetic image", "GT mask", "Prediction", "FP/FN overlay"]
    gt: set[tuple[int, int]] = set()
    pred: set[tuple[int, int]] = set()
    for y in range(120):
        for x in range(120):
            if inside_ellipse(x, y, 55, 54, 22, 15) or inside_ellipse(x, y, 78, 77, 14, 18):
                gt.add((x, y))
            if inside_ellipse(x, y, 59, 57, 24, 14) or inside_ellipse(x, y, 84, 77, 11, 15) or inside_ellipse(x, y, 34, 82, 9, 7):
                pred.add((x, y))
    tp = len(gt & pred)
    fp = len(pred - gt)
    fn = len(gt - pred)
    dice = 2 * tp / (len(gt) + len(pred))
    recall = tp / len(gt)
    burden_error = (len(pred) - len(gt)) / len(gt)
    for i, label in enumerate(labels):
        xoff = 25 + i * 250
        draw.text((xoff + 32, 24), label, fill="#17202A", font=font)
        draw.rectangle((xoff, 58, xoff + panel_w, 293), outline="#C6CCD6", width=2)
        for yy in range(rendered_case_size):
            for xx in range(rendered_case_size):
                source_x = min(source_size - 1, int(xx * source_size / rendered_case_size))
                source_y = min(source_size - 1, int(yy * source_size / rendered_case_size))
                px0 = xoff + 24 + xx
                py0 = 82 + yy
                base = int(210 - 75 * math.exp(-(((source_x - 60) ** 2 + (source_y - 64) ** 2) / 1900)))
                if i == 0:
                    img.putpixel((px0, py0), (base, base, base))
                elif i == 1:
                    img.putpixel((px0, py0), (15, 118, 110) if (source_x, source_y) in gt else (238, 242, 246))
                elif i == 2:
                    img.putpixel((px0, py0), (154, 106, 22) if (source_x, source_y) in pred else (238, 242, 246))
                else:
                    if (source_x, source_y) in gt & pred:
                        img.putpixel((px0, py0), (42, 124, 84))
                    elif (source_x, source_y) in pred - gt:
                        img.putpixel((px0, py0), (180, 58, 52))
                    elif (source_x, source_y) in gt - pred:
                        img.putpixel((px0, py0), (31, 78, 121))
                    else:
                        img.putpixel((px0, py0), (238, 242, 246))
    draw.text((55, 330), f"Dice={dice:.2f}    lesion recall={recall:.2f}    burden error={burden_error:+.1%}", fill="#17202A", font=font)
    draw.text((55, 365), f"TP pixels={tp}, FP={fp}, FN={fn}; failure mechanism: small false-positive island plus shifted lesion boundary.", fill="#606977", font=small)
    img.save(path)
    return {
        "dice": round(dice, 3),
        "lesion_recall": round(recall, 3),
        "burden_error": round(burden_error, 3),
        "tp": tp,
        "fp": fp,
        "fn": fn,
        "layout": {
            "same_synthetic_case": True,
            "source_grid_pixels": source_size,
            "rendered_case_pixels": rendered_case_size,
            "panel_pixels": panel_w,
        },
    }


def draw_result_page(slide, assets: Path, manifest: dict, refs: list[str]):
    chart_path = assets / "endpoint_ranking_chart.png"
    manifest["synthetic_endpoint_data"] = draw_endpoint_chart(chart_path)
    manifest["result_page_evidence_boundary"] = "Illustrative synthetic results - not completed validation"
    slide.shapes.add_picture(str(chart_path), inch(0.78), inch(1.42), width=inch(7.8))
    rect(slide, "Evidence boundary: illustrative synthetic results - not completed validation", 0.95, 1.21, 7.4, 0.28, "soft_gold", "gold", 8.6, True)
    rect(slide, "Interpretation\nIn this synthetic example, Calibrated wins recall; burden error is lower-is-better and lowest for Calibrated.\nBaseline wins Dice only.", 9.0, 1.62, 3.25, 1.62, "soft_teal", "teal", 11.1, True)
    rect(slide, "Meeting decision\nFreeze endpoint priority before ranking methods.", 9.0, 3.65, 3.25, 1.0, "soft_gold", "gold", 12.2, True)
    add_text(slide, reference_footer(refs, "result/interval query"), 0.95, 6.55, 11.2, 0.35, 8.4, "muted")


def draw_failure_page(slide, assets: Path, manifest: dict, refs: list[str]):
    phantom_path = assets / "synthetic_segmentation_phantom.png"
    manifest["synthetic_phantom_metrics"] = draw_phantom(phantom_path)
    manifest["phantom_overlay_legend"] = {"green": "TP/overlap", "red": "FP", "blue": "FN"}
    slide.shapes.add_picture(str(phantom_path), inch(0.72), inch(1.38), width=inch(11.85))
    add_text(slide, "Overlay legend", 9.45, 4.88, 1.2, 0.16, 8.2, "muted", True)
    swatch(slide, 9.45, 5.08, "teal")
    add_text(slide, "green = TP/overlap", 9.6, 5.06, 1.45, 0.16, 8.2, "ink", False)
    swatch(slide, 11.0, 5.08, "red")
    add_text(slide, "red = FP", 11.15, 5.06, 0.7, 0.16, 8.2, "ink", False)
    swatch(slide, 11.85, 5.08, "blue")
    add_text(slide, "blue = FN", 12.0, 5.06, 0.78, 0.16, 8.2, "ink", False)
    add_text(slide, "Scientific object: same synthetic case, aligned GT/prediction/error overlay, case metric next to visual.", 0.95, 6.35, 11.2, 0.32, 12.2, "ink", True)
    add_text(slide, reference_footer(refs, "failure-case query"), 0.95, 6.72, 11.2, 0.25, 8.4, "muted")


def draw_experiment_page(slide, assets: Path, manifest: dict, refs: list[str]):
    for i, site in enumerate(["Center A", "Center B", "Center C"]):
        rect(slide, f"{site}\nlocal image + local label\nexperimental unit: lesion-case", 0.75, 1.45 + i * 1.05, 2.75, 0.78, "soft_teal", "teal", 10.6, True)
        rect(slide, "local estimator\ncalibration score", 3.85, 1.5 + i * 1.05, 1.65, 0.65, "soft_blue", "blue", 10.3, True)
        line(slide, 3.55, 1.84 + i * 1.05, 3.82, 1.84 + i * 1.05, "purple", 1.6)
        line(slide, 5.55, 1.84 + i * 1.05, 6.12, 2.92, "purple", 1.4)
        line(slide, 5.55, 1.84 + i * 1.05, 8.48, 3.78, "muted", 1.1)
    rect(slide, "Transmitted\nlikelihood / score / summary update\n(raw images stay local)", 6.1, 2.38, 2.05, 1.25, "soft_gold", "gold", 10.7, True)
    line(slide, 8.2, 2.92, 8.48, 2.25, "purple", 1.8)
    rect(slide, "Global estimator\nsummary-sharing analysis", 8.52, 1.82, 2.05, 0.86, "soft_blue", "blue", 10.7, True)
    rect(slide, "Local-only comparator\nno shared summaries", 8.52, 3.32, 2.05, 0.86, "soft_teal", "teal", 10.7, True)
    rect(slide, "Endpoint evaluation\nlesion recall + burden error\nworst-center gap < 5%", 10.95, 2.55, 1.95, 1.15, "soft_red", "red", 10.4, True)
    line(slide, 10.62, 2.25, 10.92, 3.0, "purple", 1.8)
    line(slide, 10.62, 3.78, 10.92, 3.28, "purple", 1.8)
    add_text(slide, "Comparator path: local-only uses the same local scores without summary sharing; both branches enter the same endpoint gate.", 0.95, 5.75, 11.3, 0.42, 12.2, "ink", True)
    manifest["experiment_design_paths"] = {
        "explicit_local_only_comparator_branch": True,
        "endpoint_gate": "lesion recall + burden error + worst-center gap",
        "structural_connectors": ["summaries_to_global", "local_scores_to_local_only", "global_to_endpoint", "local_only_to_endpoint"],
    }
    add_text(slide, reference_footer(refs, "method/design query"), 0.95, 6.55, 11.2, 0.35, 8.4, "muted")


def draw_model_page(slide, assets: Path, manifest: dict, refs: list[str]):
    add_text(slide, "Target estimand", 0.95, 1.4, 2.5, 0.25, 12, "muted", True)
    rect(slide, "beta1: effect of true lesion burden T_i\non downstream outcome Y_i", 0.9, 1.75, 2.9, 1.0, "soft_teal", "teal", 11.2, True)
    rect(slide, "Observed AI phenotype\nM_i = T_i + U_i\nerror varies by center", 4.15, 1.75, 2.7, 1.0, "soft_gold", "gold", 11.2, True)
    rect(slide, "Validation subset\nobserve both T_i and M_i\nestimate calibration error", 7.25, 1.75, 2.7, 1.0, "soft_blue", "blue", 11.2, True)
    rect(slide, "Inference target\nY_i = beta0 + beta1 T_i + eps_i\nnaive M_i attenuates beta1", 10.05, 1.75, 2.55, 1.0, "soft_red", "red", 10.6, True)
    for x in [3.85, 6.9, 9.95]:
        line(slide, x, 2.25, x + 0.25, 2.25, "purple", 1.8)
    add_text(slide, "Toy calibration: E[T|M,center] = alpha_center + gamma_center M. Report beta1 after validation-calibrated correction, not a raw AI burden coefficient.", 1.0, 3.55, 11.0, 0.45, 13, "ink", True)
    rect(slide, "Evidence boundary\nThis page is a generated-variable inference mechanism,\nnot completed clinical validity evidence.", 2.5, 4.65, 8.3, 0.95, "soft_teal", "teal", 12.2, True)
    add_text(slide, reference_footer(refs, "estimator/model query"), 0.95, 6.55, 11.2, 0.35, 8.4, "muted")


SLIDES = [
    ("Endpoint Choice Changes the Method Ranking", "The method story changes when the endpoint changes.", "RESULT_FIGURE", draw_result_page),
    ("Average Dice Can Hide the Hard Case", "A case page binds image, GT, prediction, error overlay, and metric.", "FAILURE_CASE", draw_failure_page),
    ("Limited-Information Multi-Center Experiment", "Raw data stay local; shared summaries are judged by endpoint-level evidence.", "EXPERIMENT_DESIGN", draw_experiment_page),
    ("AI Phenotype Error Changes Inference", "Segmentation output is a measured variable whose error can attenuate inference.", "STATISTICAL_MODEL", draw_model_page),
]


def build_pptx(path: Path, assets: Path, manifest: dict) -> None:
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)
    blank = prs.slide_layouts[6]
    manifest["slides"] = []
    for index, (title, message, archetype, drawer) in enumerate(SLIDES, start=1):
        retrieval = retrieve_references(archetype)
        refs = list(retrieval["selected_ids"])
        slide = prs.slides.add_slide(blank)
        header(slide, index, title, message)
        drawer(slide, assets, manifest, refs)
        manifest["slides"].append({
            "slide": index,
            "title": title,
            "archetype": archetype,
            "reference_ids": refs,
            "reference_retrieval": retrieval,
            "learned_organization": retrieval["organization_lesson"],
            "reference_rationale": "References are selected by an auditable query over inspected page records with rendered-page checksums and inspection evidence; they supply page-function patterns, not copied visual assets.",
            "style_not_copied": "No whole-slide screenshot, public slide styling, private CARE figure, or clinical image is copied.",
            "what_not_copied": retrieval["what_was_not_copied"],
            "expected_scientific_objects": {
                "RESULT_FIGURE": ["endpoint-wise data", "error intervals", "method ranking", "endpoint decision"],
                "FAILURE_CASE": ["synthetic image", "GT mask", "prediction mask", "FP/FN overlay", "case metrics"],
                "EXPERIMENT_DESIGN": ["centers", "local image/label", "local estimator", "transmitted summary", "global estimator", "endpoint evaluation"],
                "STATISTICAL_MODEL": ["estimand", "observed AI phenotype", "measurement error", "validation subset", "inference target"],
            }[archetype],
        })
    prs.save(path)


def editable_slide_count(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])


def find_renderer() -> str | None:
    explicit = os.environ.get("RESEARCH_PPTX_RENDERER")
    if explicit and Path(explicit).exists():
        return explicit
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    tools = Path(".cache/tools")
    if tools.exists():
        for extracted_soffice in sorted(tools.glob("squashfs-root/opt/libreoffice*/program/soffice")):
            return str(extracted_soffice)
        extracted = tools / "squashfs-root" / "AppRun"
        if extracted.exists():
            return str(extracted)
        for candidate in sorted(tools.glob("LibreOffice*.AppImage")):
            return str(candidate)
    return None


def render_pptx(pptx_path: Path, out_dir: Path) -> dict:
    renderer = find_renderer()
    render_dir = out_dir / "rendered"
    pdf_dir = out_dir / "pdf"
    profile_dir = out_dir / "lo-profile"
    home_dir = out_dir / "lo-home"
    cache_dir = out_dir / "fontconfig-cache"
    render_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    home_dir.mkdir(parents=True, exist_ok=True)
    cache_dir.mkdir(parents=True, exist_ok=True)
    if not renderer:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "reason": "no soffice/libreoffice/RESEARCH_PPTX_RENDERER found", "png_count": 0}
    profile_uri = profile_dir.resolve().as_uri()
    env = os.environ.copy()
    env["HOME"] = str(home_dir.resolve())
    env["XDG_CACHE_HOME"] = str(cache_dir.resolve())
    cmd = [
        renderer,
        f"-env:UserInstallation={profile_uri}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir.resolve()),
        str(pptx_path.resolve()),
    ]
    result = subprocess.run(cmd, check=False, capture_output=True, text=True, timeout=120, env=env)
    pdf_candidates = sorted(pdf_dir.glob("*.pdf"))
    if result.returncode != 0 or not pdf_candidates:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "returncode": result.returncode, "stderr": result.stderr[-1000:], "png_count": 0}
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pdftoppm not found after real PPTX-to-PDF conversion", "png_count": 0}
    prefix = render_dir / "slide"
    ppm = subprocess.run([pdftoppm, "-png", "-r", "160", str(pdf_candidates[0]), str(prefix)], check=False, capture_output=True, text=True, timeout=120)
    pngs = sorted(render_dir.glob("slide-*.png"))
    return {
        "status": "ok" if ppm.returncode == 0 and pngs else "BLOCKED_REAL_PPTX_RENDER",
        "renderer": renderer,
        "pptx_to_pdf": str(pdf_candidates[0]),
        "png_count": len(pngs),
        "rendered_pngs": [str(path) for path in pngs],
        "returncode": ppm.returncode,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out-dir", type=Path, default=Path(".cache/research-group-meeting-regression"))
    args = parser.parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)
    assets = args.out_dir / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    pptx_path = args.out_dir / "research_group_meeting_regression.pptx"
    manifest_path = args.out_dir / "EVIDENCE_MANIFEST.json"
    render_status_path = args.out_dir / "RENDER_STATUS.json"
    manifest: dict = {
        "status": "GENERATED_SOURCE_ARTIFACTS_ONLY",
        "generator_may_pass": False,
        "rights_note": "No downloaded public deck, private CARE figure, whole-slide screenshot, or clinical image is copied.",
    }
    build_pptx(pptx_path, assets, manifest)
    manifest["pptx"] = str(pptx_path)
    manifest["editable_slide_count"] = editable_slide_count(pptx_path)
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    render = render_pptx(pptx_path, args.out_dir)
    render_status_path.write_text(json.dumps(render, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"pptx": str(pptx_path), "evidence_manifest": str(manifest_path), "render_status": str(render_status_path)}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
