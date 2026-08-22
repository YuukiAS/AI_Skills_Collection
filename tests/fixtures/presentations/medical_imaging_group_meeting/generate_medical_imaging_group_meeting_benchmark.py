#!/usr/bin/env python3
"""Generate a five-slide medical-imaging group-meeting benchmark deck."""

from __future__ import annotations

import argparse
import csv
import json
import math
import os
import random
import shutil
import subprocess
from pathlib import Path
from zipfile import ZipFile

os.environ.setdefault("MPLCONFIGDIR", "/tmp/ai-skills-matplotlib")
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


W, H = 13.333, 7.5
REPO_ROOT = Path(__file__).resolve().parents[4]
REFERENCE_ROOT = REPO_ROOT / "skills" / "tools" / "documents-media" / "presentations" / "shared" / "references"
REFERENCE_INDEX = REFERENCE_ROOT / "research_slide_reference_index.csv"
REFERENCE_MANIFEST = REFERENCE_ROOT / "reference_sources_manifest.json"
TASK_KEY = "017_medical_imaging_group_meeting_benchmark"
SEED = 20260822
SIZE = 192
CASES_PER_CENTER = 30

P = {
    "bg": "F4F5F2",
    "ink": "1C2430",
    "muted": "69717D",
    "line": "C7CDD3",
    "navy": "243B53",
    "teal": "0F766E",
    "green": "2A7C54",
    "red": "B0473F",
    "blue": "2F66B3",
    "gold": "9A6A16",
    "white": "FFFFFF",
    "soft": "E9ECE8",
    "panel": "FBFBF8",
}

OVERLAY_RGBA = {
    "gt": (154, 106, 22, 155),
    "tp": (42, 124, 84, 170),
    "prediction": (42, 124, 84, 150),
    "fp": (47, 102, 179, 170),
    "fn": (176, 71, 63, 190),
}

CENTER_CONFIGS = [
    {"center": "Center A", "condition": "reference contrast", "contrast": 1.00, "noise": 5.0, "bias": 0.02, "miss_small": 0.10, "fp": 0.02},
    {"center": "Center B", "condition": "low contrast + bias", "contrast": 0.78, "noise": 8.0, "bias": 0.13, "miss_small": 0.35, "fp": 0.04},
    {"center": "Center C", "condition": "high shift", "contrast": 0.58, "noise": 12.0, "bias": -0.18, "miss_small": 0.72, "fp": 0.08},
]

REFERENCE_QUERIES = {
    "MEDICAL_IMAGE_COMPARISON": {
        "intent": "Ground a segmentation task with actual image pixels, anatomy, lesion target, GT/prediction overlay, and endpoint semantics.",
        "page_functions": ["MEDICAL_IMAGE_COMPARISON", "REAL_DATA_APPLICATION", "METHOD_DIAGRAM", "MODEL_CHECK"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "cardiac MRI", "Bayesian clinical research", "hybrid resource-bound analysis"],
        "evidence_types": ["representative samples", "segmentation mask", "task overview", "metric definition", "qualitative comparison"],
        "organization_lesson": "Let the image and mask relationship be the main object before introducing aggregate metrics.",
    },
    "EXPERIMENT_DESIGN": {
        "intent": "Show center-specific appearance shift, synthetic image/GT generation, prediction, endpoint calculation, and aggregation.",
        "page_functions": ["EXPERIMENT_DESIGN", "METHOD_DIAGRAM", "SIMULATION", "MEDICAL_IMAGE_COMPARISON"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "cardiac MRI", "Bayesian workflow", "hybrid resource-bound analysis"],
        "evidence_types": ["task overview", "method mechanism", "simulation comparison", "estimator pipeline", "segmentation mask"],
        "organization_lesson": "Use one clear computation path and make image examples do part of the explanatory work.",
    },
    "RESULT_FIGURE": {
        "intent": "Compare Dice, lesion-level recall, and false-positive burden across center shift with uncertainty or variation.",
        "page_functions": ["RESULT_FIGURE", "SENSITIVITY_ANALYSIS", "REAL_DATA_APPLICATION", "MODEL_CHECK"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "cardiac MRI", "Bayesian clinical research", "Bayesian inference design"],
        "evidence_types": ["quantitative plot", "uncertainty interval", "subgroup result", "metric definition", "sensitivity analysis"],
        "organization_lesson": "Make endpoint disagreement visible in one aligned figure instead of using separate metric cards.",
    },
    "FAILURE_CASE": {
        "intent": "Show one same-case input, GT, prediction, and error overlay with directly interpretable color legend and case metrics.",
        "page_functions": ["MEDICAL_IMAGE_COMPARISON", "NEGATIVE_RESULT", "MODEL_CHECK", "REAL_DATA_APPLICATION"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "cardiac MRI", "Bayesian clinical research"],
        "evidence_types": ["qualitative comparison", "segmentation mask", "negative/fix title", "representative samples", "limitation analysis"],
        "organization_lesson": "Align same-case panels and put the error color legend close to the pixels.",
    },
    "NEGATIVE_RESULT": {
        "intent": "Show the supported failure regime and a planned validation experiment without presenting it as completed evidence.",
        "page_functions": ["NEGATIVE_RESULT", "SENSITIVITY_ANALYSIS", "NEXT_EXPERIMENT", "FINITE_SAMPLE", "RESULT_FIGURE"],
        "scientific_domain": ["medical imaging", "statistics", "biostatistics"],
        "statistical_subdomain": ["lesion segmentation", "cardiac MRI", "Bayesian clinical research", "Bayesian inference design"],
        "evidence_types": ["negative/fix title", "sensitivity analysis", "planned evidence", "subgroup result", "quantitative plot"],
        "organization_lesson": "Keep the failure evidence dominant and make the next validation question a bounded proposal.",
    },
}

TIER_PRIORITY = {
    "PRIMARY_RESEARCH_PRESENTATION": 0,
    "SECONDARY_TEACHING_REFERENCE": 1,
    "PRESENTATION_GUIDANCE": 2,
    "CANDIDATE_BACKLOG": 3,
}

REFERENCE_DESIGN_DECISIONS = {
    "MEDICAL_IMAGE_COMPARISON": {
        "adopted_design_decisions": [
            "Make the slice/overlay the largest object so the task is readable from pixels.",
            "Place endpoint and synthetic scope adjacent to the image, not in separate cards.",
            "Use direct lesion/anatomy labels instead of a detached glossary.",
        ],
        "deliberately_not_adopted": [
            "Do not copy source images, clinical cases, color styling, or institutional layout.",
            "Do not print reference IDs or retrieval traces on the audience-facing slide.",
        ],
    },
    "EXPERIMENT_DESIGN": {
        "adopted_design_decisions": [
            "Use one left-to-right data path from center shift to endpoint aggregation.",
            "Use compact image thumbnails as scientific objects inside the path.",
            "Keep connectors semantic and visible instead of decorative workflow arrows.",
        ],
        "deliberately_not_adopted": [
            "Do not show diagram QA language or diagram contracts.",
            "Do not make every operation a pastel dashboard card.",
        ],
    },
    "RESULT_FIGURE": {
        "adopted_design_decisions": [
            "Let aligned quantitative panels dominate the page.",
            "Show cross-case variation/uncertainty in the plot grammar.",
            "Annotate the endpoint disagreement directly at the high-shift condition.",
        ],
        "deliberately_not_adopted": [
            "Do not replace the result with three isolated metric tiles.",
            "Do not imply all endpoints share the same favorable direction.",
        ],
    },
    "FAILURE_CASE": {
        "adopted_design_decisions": [
            "Use same-case aligned panels so GT, prediction, and error are spatially comparable.",
            "Put TP/FP/FN color semantics in a clear projection-scale legend.",
            "Tie case metrics to the visible missed lesion region.",
        ],
        "deliberately_not_adopted": [
            "Do not crop the lesion into a tiny inset that cannot be inspected.",
            "Do not combine panels from different cases for layout convenience.",
        ],
    },
    "NEGATIVE_RESULT": {
        "adopted_design_decisions": [
            "Make the small-lesion failure curve the main evidence object.",
            "Separate completed synthetic evidence from the planned held-out-center validation.",
            "Keep the mechanism statement short and directly supported by the plotted evidence.",
        ],
        "deliberately_not_adopted": [
            "Do not build three status cards for failure, mechanism, and next work.",
            "Do not describe planned validation as a completed experiment.",
        ],
    },
}


def rgb(name: str) -> RGBColor:
    value = P[name]
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def hex_color(name: str) -> str:
    return "#" + P[name]


def inch(value: float):
    return Inches(value)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float = 13, color: str = "ink", bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def add_panel_label(slide, text: str, x: float, y: float, w: float, color: str = "ink"):
    return add_text(slide, text, x, y, w, 0.24, 10.5, color, True, PP_ALIGN.CENTER)


def add_rule(slide, x1: float, y1: float, x2: float, y2: float, color: str = "line", width: float = 1.1, dash: bool = False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dash:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return connector


def arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = "teal", width: float = 1.6):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    ln = connector.line._get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    ln.append(tail)
    return connector


def header(slide, number: int, title: str, message: str):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0), inch(0), inch(W), inch(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb("navy")
    band.line.fill.background()
    add_text(slide, title, 0.55, 0.34, 10.3, 0.45, 20.5, "ink", True)
    add_text(slide, f"{number:02d}/05", 11.65, 0.42, 0.92, 0.22, 9, "muted", False, PP_ALIGN.RIGHT)
    add_text(slide, message, 0.65, 0.88, 11.75, 0.36, 12.1, "teal", True)


def load_font(size: int):
    for path in ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/dejavu/DejaVuSans.ttf"]:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def tokens(value: str) -> set[str]:
    return {part.lower() for part in value.replace("/", " ").replace("-", " ").replace("_", " ").split() if len(part) > 2}


def load_reference_rows() -> list[dict[str, str]]:
    manifest = json.loads(REFERENCE_MANIFEST.read_text(encoding="utf-8"))
    source_tiers = {item["source_id"]: item["source_tier"] for item in manifest["candidate_sources"]}
    source_domains = {item["source_id"]: item["domain_family"] for item in manifest["candidate_sources"]}
    source_subdomains = {item["source_id"]: item["statistical_subdomain"] for item in manifest["candidate_sources"]}
    rows = []
    with REFERENCE_INDEX.open(encoding="utf-8", newline="") as fh:
        for row in csv.DictReader(fh):
            if row.get("verification_status") != "inspected":
                continue
            if not (row.get("source_file_sha256") and row.get("rendered_page_sha256")):
                continue
            row["source_tier"] = source_tiers.get(row["source_id"], "")
            row["domain_family"] = source_domains.get(row["source_id"], "")
            row["statistical_subdomain"] = source_subdomains.get(row["source_id"], "")
            rows.append(row)
    return rows


def score_reference(row: dict[str, str], query: dict[str, object]) -> tuple[int, list[str]]:
    score = 0
    reasons: list[str] = []
    if row["page_function"] in query["page_functions"]:
        score += 10
        reasons.append(f"page_function={row['page_function']} matches query")
    if row["evidence_type"] in query["evidence_types"]:
        score += 6
        reasons.append(f"evidence_type={row['evidence_type']} matches query")
    if row["domain_family"] in query["scientific_domain"]:
        score += 4
        reasons.append(f"domain={row['domain_family']} matches query")
    if row["statistical_subdomain"] in query["statistical_subdomain"]:
        score += 4
        reasons.append(f"subdomain={row['statistical_subdomain']} matches query")
    query_terms = tokens(" ".join([str(query["intent"]), " ".join(query["page_functions"]), " ".join(query["evidence_types"])]))
    row_terms = tokens(" ".join([row["scientific_object"], row["why_this_specific_page_works"], row["what_to_learn"], row["short_page_specific_observation"]]))
    overlap = sorted(query_terms & row_terms)
    if overlap:
        score += min(5, len(overlap))
        reasons.append("semantic tokens overlap: " + ", ".join(overlap[:6]))
    tier_bonus = max(0, 3 - TIER_PRIORITY.get(row["source_tier"], 3))
    score += tier_bonus
    reasons.append(f"source_tier={row['source_tier']} priority_bonus={tier_bonus}")
    return score, reasons


def retrieve_references(archetype: str) -> dict[str, object]:
    query = REFERENCE_QUERIES[archetype]
    scored = []
    for row in load_reference_rows():
        score, reasons = score_reference(row, query)
        if score > 0:
            scored.append((score, TIER_PRIORITY.get(row["source_tier"], 3), row["reference_id"], row, reasons))
    scored.sort(key=lambda item: (-item[0], item[1], item[2]))
    candidate_rows = scored[:8]
    selected_rows = candidate_rows[:5]
    if len(selected_rows) < 2:
        raise RuntimeError(f"reference retrieval for {archetype} found fewer than two inspected references")
    return {
        "query": query,
        "candidate_ids": [row["reference_id"] for _, _, _, row, _ in candidate_rows],
        "selected_ids": [row["reference_id"] for _, _, _, row, _ in selected_rows],
        "source_tiers": {row["reference_id"]: row["source_tier"] for _, _, _, row, _ in selected_rows},
        "ranking_relevance_reason": {row["reference_id"]: "; ".join(reasons) for _, _, _, row, reasons in selected_rows},
        "selected_page_lessons": {
            row["reference_id"]: {
                "primary_object": row["scientific_object"],
                "page_function": row["page_function"],
                "evidence_type": row["evidence_type"],
                "information_density": row["approximate_figure_text_ratio"],
                "annotation_style": row["short_page_specific_observation"],
                "useful_lesson": row["what_to_learn"],
                "not_copied": row["what_not_to_copy"],
            }
            for _, _, _, row, _ in selected_rows
        },
        "organization_lesson": query["organization_lesson"],
        "what_was_not_copied": "No source image, clinical/patient data, whole-slide screenshot, public deck styling, or reference visual identity was copied.",
    }


def reference_design_audit(archetype: str, retrieval: dict[str, object]) -> dict[str, object]:
    decisions = REFERENCE_DESIGN_DECISIONS[archetype]
    return {
        "selected_reference_ids": retrieval["selected_ids"],
        "page_specific_lessons": retrieval["selected_page_lessons"],
        "adopted_design_decisions": decisions["adopted_design_decisions"],
        "deliberately_not_adopted": decisions["deliberately_not_adopted"],
        "current_slide_application": retrieval["organization_lesson"],
    }


def ellipse_points(cx: float, cy: float, rx: float, ry: float) -> list[tuple[int, int]]:
    points = []
    for y in range(SIZE):
        for x in range(SIZE):
            if ((x - cx) / rx) ** 2 + ((y - cy) / ry) ** 2 <= 1:
                points.append((x, y))
    return points


def mask_from_ellipse(cx: float, cy: float, rx: float, ry: float) -> set[tuple[int, int]]:
    return set(ellipse_points(cx, cy, rx, ry))


def draw_mask_outline(draw: ImageDraw.ImageDraw, mask: set[tuple[int, int]], color: tuple[int, int, int, int], width: int = 2):
    if not mask:
        return
    xs = [x for x, _ in mask]
    ys = [y for _, y in mask]
    draw.ellipse([min(xs), min(ys), max(xs), max(ys)], outline=color, width=width)


def draw_mask_boundary(draw: ImageDraw.ImageDraw, mask: set[tuple[int, int]], color: tuple[int, int, int, int], width: int = 3):
    if not mask:
        return
    xs = [x for x, _ in mask]
    ys = [y for _, y in mask]
    pad = max(2, width)
    draw.ellipse([min(xs) - pad, min(ys) - pad, max(xs) + pad, max(ys) + pad], outline=color, width=width)


def make_case(center_index: int, case_index: int, config: dict[str, object]) -> dict[str, object]:
    rng = random.Random(SEED + center_index * 1000 + case_index)
    lesion_size = ["small", "medium", "large"][case_index % 3]
    radius = {"small": rng.uniform(3.2, 4.8), "medium": rng.uniform(6.0, 8.0), "large": rng.uniform(9.5, 12.5)}[lesion_size]
    angle = rng.uniform(-1.2, 1.5)
    cx = 96 + math.cos(angle) * rng.uniform(30, 42)
    cy = 96 + math.sin(angle) * rng.uniform(25, 36)
    outer = mask_from_ellipse(96, 96, 72, 60)
    inner = mask_from_ellipse(96, 96, 40, 32)
    myocardium = outer - inner
    gt = mask_from_ellipse(cx, cy, radius * 1.18, radius * 0.92) & myocardium
    image = Image.new("L", (SIZE, SIZE), 22)
    pix = image.load()
    contrast = float(config["contrast"])
    noise = float(config["noise"])
    bias = float(config["bias"])
    for y in range(SIZE):
        for x in range(SIZE):
            gradient = bias * (x - SIZE / 2) / (SIZE / 2) * 38
            value = 24 + gradient + rng.gauss(0, noise * 0.35)
            if (x, y) in outer:
                value = 62 + 36 * contrast + gradient + rng.gauss(0, noise)
            if (x, y) in inner:
                value = 38 + 16 * contrast + gradient + rng.gauss(0, noise * 0.8)
            if (x, y) in gt:
                value = 122 + 32 * contrast + gradient + rng.gauss(0, noise * 0.7)
            pix[x, y] = int(max(0, min(255, value)))
    image = image.filter(ImageFilter.GaussianBlur(radius=0.6))
    if lesion_size == "small":
        miss_threshold = float(config["miss_small"])
    elif lesion_size == "medium":
        miss_threshold = 0.0
    else:
        miss_threshold = 0.0
    detected = rng.random() > miss_threshold
    pred: set[tuple[int, int]] = set()
    if detected:
        if lesion_size == "small":
            shift = center_index * rng.uniform(0.8, 2.0)
            shrink = 0.88 - center_index * 0.08
        else:
            shift = center_index * rng.uniform(0.4, 1.2)
            shrink = 0.96 - center_index * 0.04
        pred = mask_from_ellipse(cx + shift, cy - shift * 0.6, radius * 1.18 * shrink, radius * 0.92 * shrink)
    if rng.random() < float(config["fp"]) + center_index * 0.05:
        fp_angle = angle + rng.uniform(0.8, 1.8)
        fp_cx = 96 + math.cos(fp_angle) * rng.uniform(32, 48)
        fp_cy = 96 + math.sin(fp_angle) * rng.uniform(26, 42)
        pred |= mask_from_ellipse(fp_cx, fp_cy, rng.uniform(3.0, 6.2), rng.uniform(2.6, 5.4))
    pred &= myocardium
    tp = len(gt & pred)
    fp = len(pred - gt)
    fn = len(gt - pred)
    dice = (2 * tp / max(1, len(gt) + len(pred))) if (gt or pred) else 1.0
    lesion_recall = 1.0 if tp / max(1, len(gt)) >= 0.18 else 0.0
    fp_burden = fp / max(1, len(myocardium))
    return {
        "case_id": f"C{center_index + 1}-{case_index:02d}",
        "center": config["center"],
        "condition": config["condition"],
        "lesion_size": lesion_size,
        "image": image,
        "gt": gt,
        "pred": pred,
        "myocardium": myocardium,
        "lesion_center": [round(cx, 2), round(cy, 2)],
        "metrics": {
            "dice": round(dice, 4),
            "lesion_recall": round(lesion_recall, 4),
            "fp_burden": round(fp_burden, 4),
            "gt_area": len(gt),
            "pred_area": len(pred),
            "tp": tp,
            "fp": fp,
            "fn": fn,
        },
    }


def summarize(values: list[float]) -> dict[str, float]:
    mean = sum(values) / len(values)
    if len(values) > 1:
        sd = math.sqrt(sum((value - mean) ** 2 for value in values) / (len(values) - 1))
    else:
        sd = 0.0
    return {"mean": round(mean, 4), "sd": round(sd, 4), "se": round(sd / math.sqrt(len(values)), 4)}


def generate_dataset(out_dir: Path) -> dict[str, object]:
    cases = []
    for center_index, config in enumerate(CENTER_CONFIGS):
        for case_index in range(CASES_PER_CENTER):
            cases.append(make_case(center_index, case_index, config))
    center_summary = []
    for config in CENTER_CONFIGS:
        subset = [case for case in cases if case["center"] == config["center"]]
        center_summary.append({
            "center": config["center"],
            "condition": config["condition"],
            "n_cases": len(subset),
            "dice": summarize([case["metrics"]["dice"] for case in subset]),
            "lesion_recall": summarize([case["metrics"]["lesion_recall"] for case in subset]),
            "fp_burden": summarize([case["metrics"]["fp_burden"] for case in subset]),
        })
    size_summary = []
    for config in CENTER_CONFIGS:
        for size in ["small", "medium", "large"]:
            subset = [case for case in cases if case["center"] == config["center"] and case["lesion_size"] == size]
            size_summary.append({
                "center": config["center"],
                "condition": config["condition"],
                "lesion_size": size,
                "n_cases": len(subset),
                "dice": summarize([case["metrics"]["dice"] for case in subset]),
                "lesion_recall": summarize([case["metrics"]["lesion_recall"] for case in subset]),
                "fp_burden": summarize([case["metrics"]["fp_burden"] for case in subset]),
            })
    failure_case = min(
        [case for case in cases if case["lesion_size"] == "small" and case["center"] == "Center C"],
        key=lambda case: (case["metrics"]["lesion_recall"], case["metrics"]["dice"]),
    )
    summary = {
        "seed": SEED,
        "cases_per_center": CASES_PER_CENTER,
        "centers": [{"center": c["center"], "condition": c["condition"], "contrast": c["contrast"], "noise": c["noise"], "bias": c["bias"]} for c in CENTER_CONFIGS],
        "endpoints": ["Dice overlap", "lesion-level recall", "false-positive burden"],
        "center_summary": center_summary,
        "lesion_size_summary": size_summary,
        "failure_case_id": failure_case["case_id"],
        "negative_result": {
            "condition": "Center C high shift, small lesions",
            "small_lesion_recall": next(item["lesion_recall"]["mean"] for item in size_summary if item["center"] == "Center C" and item["lesion_size"] == "small"),
            "center_c_dice": next(item["dice"]["mean"] for item in center_summary if item["center"] == "Center C"),
            "planned_validation": "planned held-out-center lesion-size stratification with threshold/calibration sensitivity",
        },
    }
    serializable_cases = []
    for case in cases:
        serializable_cases.append({
            "case_id": case["case_id"],
            "center": case["center"],
            "condition": case["condition"],
            "lesion_size": case["lesion_size"],
            "lesion_center": case["lesion_center"],
            "metrics": case["metrics"],
        })
    sim_dir = out_dir / "simulation"
    sim_dir.mkdir(parents=True, exist_ok=True)
    (sim_dir / "case_metrics.json").write_text(json.dumps(serializable_cases, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (sim_dir / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    with (sim_dir / "center_summary.csv").open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["center", "condition", "n_cases", "dice_mean", "dice_se", "lesion_recall_mean", "lesion_recall_se", "fp_burden_mean", "fp_burden_se"])
        for row in center_summary:
            writer.writerow([row["center"], row["condition"], row["n_cases"], row["dice"]["mean"], row["dice"]["se"], row["lesion_recall"]["mean"], row["lesion_recall"]["se"], row["fp_burden"]["mean"], row["fp_burden"]["se"]])
    return {"cases": cases, "summary": summary}


def rgba_from_mask(base: Image.Image, gt: set[tuple[int, int]], pred: set[tuple[int, int]], mode: str) -> Image.Image:
    rgb = base.convert("RGB")
    overlay = Image.new("RGBA", rgb.size, (0, 0, 0, 0))
    pix = overlay.load()
    if mode == "gt_pred":
        for x, y in gt:
            pix[x, y] = OVERLAY_RGBA["gt"]
        for x, y in pred:
            pix[x, y] = OVERLAY_RGBA["prediction"]
    elif mode == "gt":
        for x, y in gt:
            pix[x, y] = OVERLAY_RGBA["gt"]
    elif mode == "pred":
        for x, y in pred:
            pix[x, y] = OVERLAY_RGBA["prediction"]
    elif mode == "error":
        for x, y in gt & pred:
            pix[x, y] = OVERLAY_RGBA["tp"]
        for x, y in pred - gt:
            pix[x, y] = OVERLAY_RGBA["fp"]
        for x, y in gt - pred:
            pix[x, y] = OVERLAY_RGBA["fn"]
    draw = ImageDraw.Draw(overlay)
    if mode in {"gt_pred", "gt"}:
        draw_mask_boundary(draw, gt, OVERLAY_RGBA["gt"], width=3)
    if mode in {"gt_pred", "pred"}:
        draw_mask_boundary(draw, pred, OVERLAY_RGBA["prediction"], width=3)
    if mode == "error":
        draw_mask_boundary(draw, gt & pred, OVERLAY_RGBA["tp"], width=3)
        draw_mask_boundary(draw, pred - gt, OVERLAY_RGBA["fp"], width=3)
        draw_mask_boundary(draw, gt - pred, OVERLAY_RGBA["fn"], width=4)
    return Image.alpha_composite(rgb.convert("RGBA"), overlay)


def annotate_image(
    image: Image.Image,
    label: str,
    lesion_xy: list[float] | None = None,
    marker_color: tuple[int, int, int, int] = OVERLAY_RGBA["gt"],
) -> Image.Image:
    canvas = image.convert("RGBA").resize((420, 420), Image.Resampling.BICUBIC)
    draw = ImageDraw.Draw(canvas)
    font = load_font(18)
    draw.rounded_rectangle([10, 10, 214, 44], radius=5, fill=(28, 36, 48, 210))
    draw.text((20, 18), label, fill=(255, 255, 255, 255), font=font)
    if lesion_xy:
        x, y = lesion_xy
        sx, sy = x / SIZE * 420, y / SIZE * 420
        draw.line([sx - 42, sy - 34, sx - 8, sy - 8], fill=marker_color, width=3)
        draw.ellipse([sx - 7, sy - 7, sx + 7, sy + 7], outline=marker_color, width=3)
    return canvas


def save_case_assets(assets: Path, data: dict[str, object]) -> dict[str, str]:
    assets.mkdir(parents=True, exist_ok=True)
    cases = data["cases"]
    slide1_case = next(case for case in cases if case["case_id"] == "C2-04")
    failure_case = next(case for case in cases if case["case_id"] == data["summary"]["failure_case_id"])
    paths: dict[str, str] = {}
    center_examples = {
        config["center"]: next(case for case in cases if case["center"] == config["center"] and case["case_id"].endswith("-04"))
        for config in CENTER_CONFIGS
    }
    image_defs = {
        "slide1_input": annotate_image(slide1_case["image"], "Synthetic cardiac-MR-like slice", slide1_case["lesion_center"]),
        "slide1_overlay": annotate_image(rgba_from_mask(slide1_case["image"], slide1_case["gt"], slide1_case["pred"], "gt_pred"), "GT + prediction"),
        "failure_input": annotate_image(failure_case["image"], "Input", failure_case["lesion_center"]),
        "failure_gt": annotate_image(rgba_from_mask(failure_case["image"], failure_case["gt"], failure_case["pred"], "gt"), "GT lesion"),
        "failure_pred": annotate_image(rgba_from_mask(failure_case["image"], failure_case["gt"], failure_case["pred"], "pred"), "Prediction"),
        "failure_error": annotate_image(rgba_from_mask(failure_case["image"], failure_case["gt"], failure_case["pred"], "error"), "TP / FP / FN"),
        "center_a_shift": annotate_image(rgba_from_mask(center_examples["Center A"]["image"], center_examples["Center A"]["gt"], center_examples["Center A"]["pred"], "gt_pred"), "A reference"),
        "center_b_shift": annotate_image(rgba_from_mask(center_examples["Center B"]["image"], center_examples["Center B"]["gt"], center_examples["Center B"]["pred"], "gt_pred"), "B low contrast"),
        "center_c_shift": annotate_image(rgba_from_mask(center_examples["Center C"]["image"], center_examples["Center C"]["gt"], center_examples["Center C"]["pred"], "gt_pred"), "C high shift"),
    }
    for name, image in image_defs.items():
        path = assets / f"{name}.png"
        image.save(path)
        paths[name] = str(path)
    return paths


def plot_main_results(path: Path, summary: dict[str, object]) -> None:
    rows = summary["center_summary"]
    size_rows = summary["lesion_size_summary"]
    labels = [row["center"].replace("Center ", "") for row in rows]
    x = list(range(len(rows)))
    small_recall_rows = []
    for row in rows:
        small = next(item for item in size_rows if item["center"] == row["center"] and item["lesion_size"] == "small")
        small_recall_rows.append({
            "center": row["center"],
            "lesion_recall": small["lesion_recall"],
        })
    fp_max = max(row["fp_burden"]["mean"] + row["fp_burden"]["se"] * 1.96 for row in rows)
    metrics = [
        ("Dice overlap", rows, "dice", "higher is better", 0.0, 1.0, "teal"),
        ("Small-lesion recall", small_recall_rows, "lesion_recall", "higher is better", 0.0, 1.0, "red"),
        ("FP burden", rows, "fp_burden", "lower is better", 0.0, max(0.012, fp_max * 1.55), "blue"),
    ]
    fig, axes = plt.subplots(1, 3, figsize=(11.5, 3.6), dpi=220)
    fig.patch.set_facecolor(hex_color("panel"))
    for ax, (title, metric_rows, key, subtitle, ymin, ymax, color) in zip(axes, metrics, strict=True):
        means = [row[key]["mean"] for row in metric_rows]
        ses = [row[key]["se"] * 1.96 for row in metric_rows]
        ax.errorbar(x, means, yerr=ses, fmt="o", color=hex_color(color), ecolor=hex_color(color), elinewidth=2.0, capsize=5, markersize=7)
        ax.plot(x, means, color=hex_color(color), linewidth=2.0, alpha=0.7)
        ax.set_title(title, fontsize=13, weight="bold", color=hex_color("ink"))
        ax.text(0.02, 0.93, subtitle, transform=ax.transAxes, fontsize=9, color=hex_color("muted"))
        ax.set_xticks(x, labels)
        ax.set_ylim(ymin, ymax)
        ax.grid(axis="y", color="#D8DDE0", linewidth=0.8)
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.tick_params(labelsize=10, colors=hex_color("ink"))
    axes[1].annotate(
        "small lesions missed\nunder high shift",
        xy=(2, small_recall_rows[2]["lesion_recall"]["mean"]),
        xytext=(1.12, 0.28),
        arrowprops={"arrowstyle": "->", "color": hex_color("red"), "lw": 1.4},
        fontsize=9.5,
        color=hex_color("red"),
    )
    fig.text(0.04, 0.02, "Mean across deterministic synthetic cases; vertical intervals show approximately 95% cross-case uncertainty.", fontsize=9, color=hex_color("muted"))
    fig.tight_layout(rect=(0, 0.08, 1, 1))
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def plot_size_results(path: Path, summary: dict[str, object]) -> None:
    rows = summary["lesion_size_summary"]
    sizes = ["small", "medium", "large"]
    centers = [config["center"] for config in CENTER_CONFIGS]
    fig, ax = plt.subplots(figsize=(9.6, 4.3), dpi=220)
    fig.patch.set_facecolor(hex_color("panel"))
    offsets = [-0.22, 0.0, 0.22]
    colors = [hex_color("teal"), hex_color("gold"), hex_color("red")]
    for offset, center, color in zip(offsets, centers, colors, strict=True):
        means = [next(row["lesion_recall"]["mean"] for row in rows if row["center"] == center and row["lesion_size"] == size) for size in sizes]
        ses = [next(row["lesion_recall"]["se"] * 1.96 for row in rows if row["center"] == center and row["lesion_size"] == size) for size in sizes]
        xs = [i + offset for i in range(len(sizes))]
        ax.errorbar(xs, means, yerr=ses, fmt="o", color=color, ecolor=color, elinewidth=2, capsize=5, markersize=7, label=center)
        ax.plot(xs, means, color=color, linewidth=1.8, alpha=0.65)
    ax.set_xticks(range(len(sizes)), ["small lesions", "medium", "large"])
    ax.set_ylim(-0.04, 1.04)
    ax.set_ylabel("lesion-level recall", fontsize=12, color=hex_color("ink"))
    ax.grid(axis="y", color="#D8DDE0", linewidth=0.8)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(labelsize=11, colors=hex_color("ink"))
    ax.legend(loc="upper left", frameon=False, fontsize=10)
    c_row = next(row for row in rows if row["center"] == "Center C" and row["lesion_size"] == "small")
    ax.annotate(
        f"Center C small\nrecall = {c_row['lesion_recall']['mean']:.2f}",
        xy=(0 + 0.22, c_row["lesion_recall"]["mean"]),
        xytext=(0.52, 0.22),
        arrowprops={"arrowstyle": "->", "color": hex_color("red"), "lw": 1.5},
        color=hex_color("red"),
        fontsize=12,
        weight="bold",
    )
    fig.text(0.03, 0.02, "Completed evidence: deterministic synthetic phantoms only; next validation is planned, not completed.", fontsize=9.5, color=hex_color("muted"))
    fig.tight_layout(rect=(0, 0.08, 1, 1))
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_assets(assets: Path, data: dict[str, object]) -> dict[str, str]:
    paths = save_case_assets(assets, data)
    main_plot = assets / "slide3_endpoint_disagreement.png"
    size_plot = assets / "slide5_lesion_size_recall.png"
    plot_main_results(main_plot, data["summary"])
    plot_size_results(size_plot, data["summary"])
    paths["main_plot"] = str(main_plot)
    paths["size_plot"] = str(size_plot)
    return paths


def add_image(slide, path: str, x: float, y: float, w: float, h: float | None = None):
    kwargs = {"width": inch(w)}
    if h is not None:
        kwargs["height"] = inch(h)
    return slide.shapes.add_picture(path, inch(x), inch(y), **kwargs)


def add_legend(slide, x: float, y: float):
    entries = [("GT", "gold"), ("Pred / TP", "green"), ("FP", "blue"), ("FN", "red")]
    for i, (label, color) in enumerate(entries):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x + i * 1.25), inch(y), inch(0.16), inch(0.16))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        add_text(slide, label, x + i * 1.25 + 0.22, y - 0.02, 0.92, 0.22, 8.5, "muted")


def draw_slide1(slide, assets: dict[str, str], manifest: dict, refs: list[str], data: dict[str, object]):
    add_image(slide, assets["slide1_input"], 0.70, 1.42, 4.20, 4.20)
    add_image(slide, assets["slide1_overlay"], 4.95, 1.42, 4.20, 4.20)
    add_text(slide, "Evaluation unit: lesion detection + mask overlap", 9.55, 1.50, 2.85, 0.78, 15, "ink", True)
    add_text(slide, "Dice summarizes mask overlap, but a missed small lesion is counted at the lesion level.", 9.55, 2.44, 2.75, 0.76, 12.2, "muted")
    add_text(slide, "Endpoints: Dice overlap, lesion recall, false-positive burden", 9.55, 3.52, 2.75, 0.58, 12.2, "teal", True)
    add_text(slide, "Synthetic cardiac-MR-like phantoms; not clinical validation.", 9.55, 4.62, 2.75, 0.58, 11.4, "muted")
    add_rule(slide, 0.7, 5.95, 9.15, 5.95, "line")
    add_legend(slide, 1.05, 6.15)
    manifest.setdefault("audience_text_by_archetype", {})["MEDICAL_IMAGE_COMPARISON"] = [
        "Evaluation unit: lesion detection + mask overlap",
        "Dice summarizes mask overlap, but a missed small lesion is counted at the lesion level.",
        "Endpoints: Dice overlap, lesion recall, false-positive burden",
        "Synthetic cardiac-MR-like phantoms; not clinical validation.",
    ]


def node(slide, text: str, x: float, y: float, w: float, h: float, fill: str = "panel"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("line")
    shape.line.width = Pt(1)
    add_text(slide, text, x + 0.12, y + 0.12, w - 0.24, h - 0.20, 11.2, "ink", True, PP_ALIGN.CENTER)
    return shape


def draw_slide2(slide, assets: dict[str, str], manifest: dict, refs: list[str], data: dict[str, object]):
    center_assets = [("Center A", "center_a_shift"), ("Center B", "center_b_shift"), ("Center C", "center_c_shift")]
    for y, (label, key) in zip([1.34, 2.76, 4.18], center_assets, strict=True):
        add_image(slide, assets[key], 0.68, y, 1.18, 1.18)
        condition = next(config["condition"] for config in CENTER_CONFIGS if config["center"] == label)
        add_text(slide, label, 1.96, y + 0.20, 1.02, 0.24, 10.0, "ink", True)
        add_text(slide, condition, 1.96, y + 0.53, 1.12, 0.32, 8.7, "muted")
    add_text(slide, "appearance shift", 0.75, 5.56, 2.06, 0.28, 10.2, "teal", True, PP_ALIGN.CENTER)
    add_image(slide, assets["slide1_overlay"], 3.36, 2.00, 1.42, 1.42)
    add_panel_label(slide, "image + GT", 3.36, 3.54, 1.42)
    add_image(slide, assets["failure_pred"], 5.42, 2.00, 1.42, 1.42)
    add_panel_label(slide, "prediction", 5.42, 3.54, 1.42)
    node(slide, "case metrics\nDice, recall, FP", 7.56, 2.15, 1.66, 0.90)
    node(slide, "center summary\n+ lesion strata", 9.94, 2.15, 1.78, 0.90)
    arrow(slide, 2.92, 2.55, 3.36, 2.55)
    arrow(slide, 4.78, 2.55, 5.42, 2.55)
    arrow(slide, 6.84, 2.55, 7.56, 2.55)
    arrow(slide, 9.22, 2.55, 9.94, 2.55)
    add_text(slide, "The same synthetic case carries image pixels, lesion mask, prediction and metrics before center-level averaging.", 3.26, 4.20, 7.95, 0.44, 13.0, "teal", True)
    add_text(slide, "The stress test asks whether center-level average Dice hides lesion-level failures under appearance shift.", 3.26, 4.85, 7.95, 0.48, 12.0, "muted")
    add_rule(slide, 3.20, 5.58, 11.2, 5.58, "line", dash=True)
    add_text(slide, "Aggregation is descriptive: no clinical validation claim is made from these phantoms.", 3.45, 5.82, 7.50, 0.40, 11.0, "muted", False, PP_ALIGN.CENTER)
    manifest["experiment_design_diagram"] = {
        "reading_direction": "left_to_right",
        "edge_crossing": "none",
        "arrowheads": ["tailEnd"],
        "structural_connectors": ["center_to_image_gt", "image_gt_to_prediction", "prediction_to_case_metrics", "case_metrics_to_center_summary"],
    }
    manifest.setdefault("audience_text_by_archetype", {})["EXPERIMENT_DESIGN"] = [
        "3 center conditions",
        "contrast, noise, bias",
        "Center A reference contrast",
        "Center B low contrast + bias",
        "Center C high shift",
        "synthetic image + GT lesion",
        "deterministic prediction",
        "case metrics Dice, recall, FP",
        "center summary + lesion strata",
        "Case-level metrics stay attached to the same image, GT and prediction before any center-level average.",
        "Aggregation is descriptive: no clinical validation claim is made from these phantoms.",
    ]


def draw_slide3(slide, assets: dict[str, str], manifest: dict, refs: list[str], data: dict[str, object]):
    add_image(slide, assets["main_plot"], 0.72, 1.30, 11.70, 4.25)
    c = data["summary"]["negative_result"]
    add_text(slide, f"Center C average Dice = {c['center_c_dice']:.2f}, but small-lesion recall drops to {c['small_lesion_recall']:.2f}.", 0.95, 5.84, 7.60, 0.36, 13.0, "ink", True)
    add_text(slide, "FP burden is lower-is-better; Dice and recall are higher-is-better.", 8.95, 5.76, 3.05, 0.50, 10.8, "muted")
    manifest.setdefault("audience_text_by_archetype", {})["RESULT_FIGURE"] = [
        "Center C average Dice remains moderate, but small-lesion recall drops under high shift.",
        "Mean across deterministic synthetic cases; vertical intervals show approximately 95% cross-case uncertainty.",
        "FP burden is lower-is-better; Dice and recall are higher-is-better.",
    ]


def draw_slide4(slide, assets: dict[str, str], manifest: dict, refs: list[str], data: dict[str, object]):
    xs = [0.58, 3.72, 6.86, 10.00]
    labels = ["Input", "GT lesion", "Prediction", "TP / FP / FN"]
    keys = ["failure_input", "failure_gt", "failure_pred", "failure_error"]
    for x, key, label in zip(xs, keys, labels, strict=True):
        add_image(slide, assets[key], x, 1.48, 2.72, 2.72)
        add_panel_label(slide, label, x, 4.30, 2.72)
    add_legend(slide, 0.76, 4.88)
    failure = next(case for case in data["cases"] if case["case_id"] == data["summary"]["failure_case_id"])
    metrics = failure["metrics"]
    add_text(slide, f"Case {failure['case_id']} ({failure['center']}, {failure['lesion_size']} lesion): Dice {metrics['dice']:.2f}; lesion recall {metrics['lesion_recall']:.0f}; FN pixels {metrics['fn']}.", 0.78, 5.56, 8.65, 0.42, 13.0, "ink", True)
    add_text(slide, "The red FN region shows why an average overlap metric is not enough for this stress condition.", 0.78, 6.08, 8.65, 0.34, 11.2, "muted")
    add_text(slide, "All four panels are the same synthetic slice geometry.", 9.65, 5.56, 2.55, 0.42, 10.8, "muted")
    manifest["failure_case"] = {
        "case_id": failure["case_id"],
        "center": failure["center"],
        "lesion_size": failure["lesion_size"],
        "same_slice_geometry": True,
        "metrics": metrics,
        "panels": labels,
    }
    manifest.setdefault("audience_text_by_archetype", {})["FAILURE_CASE"] = [
        "Input",
        "GT lesion",
        "Prediction",
        "TP / FP / FN",
        f"Case {failure['case_id']} ({failure['center']}, {failure['lesion_size']} lesion): Dice {metrics['dice']:.2f}; lesion recall {metrics['lesion_recall']:.0f}; FN pixels {metrics['fn']}.",
        "All four panels are the same synthetic slice geometry.",
    ]


def draw_slide5(slide, assets: dict[str, str], manifest: dict, refs: list[str], data: dict[str, object]):
    add_image(slide, assets["size_plot"], 0.72, 1.35, 8.70, 4.30)
    add_text(slide, "Completed evidence", 9.72, 1.50, 2.45, 0.30, 13.0, "ink", True)
    add_text(slide, "Small lesions in the high-shift center are the supported failure regime in this deterministic phantom benchmark.", 9.72, 1.92, 2.45, 0.86, 11.5, "muted")
    add_text(slide, "Planned validation", 9.72, 3.28, 2.45, 0.30, 13.0, "teal", True)
    add_text(slide, "Hold out center C, stratify by lesion size, and sweep prediction threshold/calibration before comparing methods.", 9.72, 3.70, 2.45, 0.94, 11.5, "muted")
    add_text(slide, "Synthetic cardiac-MR-like phantoms; not a patient validation endpoint.", 9.72, 5.34, 2.45, 0.46, 10.5, "muted")
    manifest.setdefault("audience_text_by_archetype", {})["NEGATIVE_RESULT"] = [
        "Completed evidence",
        "Small lesions in the high-shift center are the supported failure regime in this deterministic phantom benchmark.",
        "Planned validation",
        "Hold out center C, stratify by lesion size, and sweep prediction threshold/calibration before comparing methods.",
        "Synthetic cardiac-MR-like phantoms; not a patient validation endpoint.",
    ]


SLIDES = [
    ("A small lesion is the unit of failure", "Synthetic cardiac-MR-like segmentation task; not clinical validation", "MEDICAL_IMAGE_COMPARISON", draw_slide1),
    ("The stress test follows one image-to-endpoint path", "Center appearance shift propagates through prediction before aggregation", "EXPERIMENT_DESIGN", draw_slide2),
    ("Average Dice can hide a lesion-level failure", "Endpoint disagreement appears under the high-shift center", "RESULT_FIGURE", draw_slide3),
    ("The failure is visible in the same-case overlay", "A missed small lesion drives the negative endpoint", "FAILURE_CASE", draw_slide4),
    ("Validate the small-lesion shift explicitly next", "Completed synthetic evidence motivates a planned held-out-center experiment", "NEGATIVE_RESULT", draw_slide5),
]


INTERNAL_LEAK_PATTERNS = [
    "RRL-",
    "Reference retrieval",
    "EVIDENCE_MANIFEST",
    "Diagram contract",
    "style not copied",
    "repo path",
    "run ID",
    "implementation commit",
    "review target",
]

ANTI_META_PHRASES = [
    "Reading target",
    "Observed in this synthetic run",
    "Evidence boundary",
    "Key takeaway",
    "Role in the deck",
    "This slide",
    "This page",
]

SOURCE_LIKE_PATTERNS = ["mask_", "pred_", "gt_", "case_id =", "Dice =", "FP_BURDEN"]


def build_pptx(path: Path, assets: dict[str, str], manifest: dict, data: dict[str, object]) -> None:
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)
    blank = prs.slide_layouts[6]
    manifest["slides"] = []
    manifest["reference_design_audit"] = {}
    for index, (title, message, archetype, drawer) in enumerate(SLIDES, start=1):
        retrieval = retrieve_references(archetype)
        refs = list(retrieval["selected_ids"])
        slide = prs.slides.add_slide(blank)
        header(slide, index, title, message)
        drawer(slide, assets, manifest, refs, data)
        audit = reference_design_audit(archetype, retrieval)
        manifest["reference_design_audit"][f"slide_{index}"] = audit
        manifest["slides"].append({
            "slide": index,
            "title": title,
            "archetype": archetype,
            "reference_ids": refs,
            "reference_retrieval": retrieval,
            "reference_design_audit": audit,
            "audience_text": manifest.get("audience_text_by_archetype", {}).get(archetype, []),
            "learned_organization": retrieval["organization_lesson"],
            "reference_rationale": "References are selected by auditable semantic query over inspected page records; they inform organization only.",
            "style_not_copied": "No source medical image, public slide styling, private clinical figure, patient image, or whole-slide screenshot is copied.",
            "what_not_copied": retrieval["what_was_not_copied"],
            "evidence_status": "deterministic_synthetic_medical_imaging_generated",
            "expected_scientific_objects": {
                "MEDICAL_IMAGE_COMPARISON": ["synthetic cardiac-MR-like slice", "GT lesion mask", "prediction overlay", "endpoint names", "synthetic-only qualifier"],
                "EXPERIMENT_DESIGN": ["center appearance shifts", "synthetic image and GT generation", "prediction step", "case metrics", "center aggregation"],
                "RESULT_FIGURE": ["Dice panel", "lesion recall panel", "FP burden panel", "uncertainty intervals", "endpoint disagreement annotation"],
                "FAILURE_CASE": ["same-case input", "GT panel", "prediction panel", "TP/FP/FN overlay", "case metrics"],
                "NEGATIVE_RESULT": ["lesion-size recall plot", "failure regime annotation", "completed evidence boundary", "planned validation experiment"],
            }[archetype],
        })
    prs.save(path)


def pptx_slide_text(pptx_path: Path) -> str:
    with ZipFile(pptx_path) as zf:
        slide_names = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        return "\n".join(zf.read(name).decode("utf-8", errors="ignore") for name in slide_names)


def run_quality_gates(pptx_path: Path, manifest: dict[str, object], data: dict[str, object]) -> dict[str, object]:
    slide_xml = pptx_slide_text(pptx_path)
    audience_text = "\n".join(
        "\n".join(slide.get("audience_text", []))
        for slide in manifest.get("slides", [])
        if isinstance(slide, dict)
    )
    failures: list[str] = []
    for pattern in INTERNAL_LEAK_PATTERNS:
        if pattern in slide_xml or pattern in audience_text:
            failures.append(f"audience-facing internal leak: {pattern}")
    for phrase in ANTI_META_PHRASES:
        if phrase in slide_xml or phrase in audience_text:
            failures.append(f"audience-facing meta language: {phrase}")
    for pattern in SOURCE_LIKE_PATTERNS:
        if pattern in audience_text:
            failures.append(f"source-like medical-imaging string leaked into audience text: {pattern}")
    if manifest.get("editable_slide_count") != 5:
        failures.append("editable deck must have five slides")
    audit = manifest.get("reference_design_audit", {})
    if not isinstance(audit, dict) or len(audit) != 5:
        failures.append("reference_design_audit must cover five slides")
    for slide_key, item in audit.items() if isinstance(audit, dict) else []:
        if len(item.get("selected_reference_ids", [])) < 2:
            failures.append(f"{slide_key} has fewer than two inspected reference IDs")
        if not item.get("adopted_design_decisions"):
            failures.append(f"{slide_key} lacks adopted design decisions")
    if manifest.get("experiment_design_diagram", {}).get("edge_crossing") != "none":
        failures.append("experiment design diagram must record no crossing")
    if manifest.get("failure_case", {}).get("same_slice_geometry") is not True:
        failures.append("failure case panels must be same slice geometry")
    negative = data["summary"]["negative_result"]
    if not (negative["center_c_dice"] > 0.40 and negative["small_lesion_recall"] < 0.45):
        failures.append("deterministic result does not support average-Dice versus small-lesion-recall disagreement")
    return {
        "status": "PASS" if not failures else "FAIL",
        "failures": failures,
        "checked_gates": [
            "audience_facing_internal_leak",
            "anti_meta_language",
            "same_case_failure_panels",
            "endpoint_disagreement_supported_by_metrics",
            "reference_design_audit",
        ],
    }


def editable_slide_count(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])


def find_renderer() -> str | None:
    explicit = os.environ.get("RESEARCH_PPTX_RENDERER")
    if explicit and Path(explicit).exists():
        return explicit
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    tools = Path(".cache/tools")
    if tools.exists():
        for extracted_soffice in sorted(tools.glob("squashfs-root/opt/libreoffice*/program/soffice")):
            return str(extracted_soffice)
        extracted = tools / "squashfs-root" / "AppRun"
        if extracted.exists():
            return str(extracted)
        for candidate in sorted(tools.glob("LibreOffice*.AppImage")):
            return str(candidate)
    return None


def render_pptx(pptx_path: Path, out_dir: Path) -> dict:
    renderer = find_renderer()
    render_dir = out_dir / "rendered"
    pdf_dir = out_dir / "pdf"
    profile_dir = out_dir / "lo-profile"
    home_dir = out_dir / "lo-home"
    cache_dir = out_dir / "fontconfig-cache"
    render_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    home_dir.mkdir(parents=True, exist_ok=True)
    cache_dir.mkdir(parents=True, exist_ok=True)
    if not renderer:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "reason": "no soffice/libreoffice/RESEARCH_PPTX_RENDERER found", "png_count": 0}
    profile_uri = profile_dir.resolve().as_uri()
    env = os.environ.copy()
    env["HOME"] = str(home_dir.resolve())
    env["XDG_CACHE_HOME"] = str(cache_dir.resolve())
    cmd = [
        renderer,
        f"-env:UserInstallation={profile_uri}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir.resolve()),
        str(pptx_path.resolve()),
    ]
    result = subprocess.run(cmd, check=False, capture_output=True, text=True, timeout=120, env=env)
    pdf_candidates = sorted(pdf_dir.glob("*.pdf"))
    if result.returncode != 0 or not pdf_candidates:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "returncode": result.returncode, "stderr": result.stderr[-1000:], "png_count": 0}
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return {"status": "BLOCKED_REAL_PPTX_RENDER", "renderer": renderer, "reason": "pdftoppm not found after real PPTX-to-PDF conversion", "png_count": 0}
    prefix = render_dir / "slide"
    ppm = subprocess.run([pdftoppm, "-png", "-r", "160", str(pdf_candidates[0]), str(prefix)], check=False, capture_output=True, text=True, timeout=120)
    pngs = sorted(render_dir.glob("slide-*.png"))
    if ppm.returncode == 0 and len(pngs) == 5:
        expected_dir = out_dir / "expected_render"
        expected_dir.mkdir(parents=True, exist_ok=True)
        for png in pngs:
            shutil.copy2(png, expected_dir / png.name)
    for transient in (profile_dir, home_dir, cache_dir):
        shutil.rmtree(transient, ignore_errors=True)
    return {
        "status": "ok" if ppm.returncode == 0 and len(pngs) == 5 else "BLOCKED_REAL_PPTX_RENDER",
        "renderer": renderer,
        "pptx_to_pdf": str(pdf_candidates[0]),
        "png_count": len(pngs),
        "rendered_pngs": [str(path) for path in pngs],
        "returncode": ppm.returncode,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out-dir", type=Path, default=Path(".cache/medical-imaging-group-meeting-benchmark"))
    args = parser.parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)
    assets_dir = args.out_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    data = generate_dataset(args.out_dir)
    assets = build_assets(assets_dir, data)
    pptx_path = args.out_dir / "medical_imaging_group_meeting_benchmark.pptx"
    manifest_path = args.out_dir / "EVIDENCE_MANIFEST.json"
    reference_audit_path = args.out_dir / "reference_design_audit.json"
    render_status_path = args.out_dir / "RENDER_STATUS.json"
    manifest: dict = {
        "status": "GENERATED_SOURCE_ARTIFACTS_ONLY",
        "generator_may_pass": False,
        "task_key": TASK_KEY,
        "rights_note": "Synthetic cardiac-MR-like phantoms only; no public slide image, private clinical figure, real patient image, or hospital data is copied.",
        "benchmark_story": "Average Dice can look moderate while high-shift small-lesion recall fails; medical-imaging slides must connect image pixels, masks, endpoints, failure mechanism, and planned validation.",
        "synthetic_dataset_summary": data["summary"],
        "asset_paths": assets,
    }
    build_pptx(pptx_path, assets, manifest, data)
    manifest["pptx"] = str(pptx_path)
    manifest["editable_slide_count"] = editable_slide_count(pptx_path)
    manifest["deterministic_quality_gates"] = run_quality_gates(pptx_path, manifest, data)
    if manifest["deterministic_quality_gates"]["status"] != "PASS":
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        raise SystemExit("deterministic presentation QA failed: " + "; ".join(manifest["deterministic_quality_gates"]["failures"]))
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reference_audit_path.write_text(json.dumps(manifest["reference_design_audit"], ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    render = render_pptx(pptx_path, args.out_dir)
    render_status_path.write_text(json.dumps(render, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"pptx": str(pptx_path), "evidence_manifest": str(manifest_path), "render_status": str(render_status_path)}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
